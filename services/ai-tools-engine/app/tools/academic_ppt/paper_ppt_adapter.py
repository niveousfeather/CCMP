from __future__ import annotations

import asyncio
import hashlib
import json
import math
import os
import posixpath
import re
import shutil
import sys
import textwrap
import zipfile
from contextlib import contextmanager
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from app.core.checkpoints import write_checkpoint
from app.core.config import get_settings
from app.core.diagnostics import format_missing_dependency_message, get_academic_ppt_diagnostics
from app.core.errors import sanitize_message
from app.core.files import safe_task_child
from app.core.model_bridge import NexusModelBridgeProvider
from app.core.search_bridge import run_academic_search_bridge
from app.tools.academic_ppt.logs import append_log

BUILTIN_TEMPLATE_ID = "school_academic_report"
BUILTIN_TEMPLATE_DIR = Path(__file__).resolve().parent / "templates" / "builtin" / BUILTIN_TEMPLATE_ID
BUILTIN_TEMPLATE_PPTX = BUILTIN_TEMPLATE_DIR / "template.pptx"
BUILTIN_TEMPLATE_JSON = BUILTIN_TEMPLATE_DIR / "template.json"
BUILTIN_TEMPLATE_COVER_BACKGROUND = BUILTIN_TEMPLATE_DIR / "cover-background.jpg"
BUILTIN_TEMPLATE_LOGO_MEDIA = "ppt/media/image9.png"
BUILTIN_TEMPLATE_ROLES = {
    "cover",
    "cover_alt",
    "toc",
    "section",
    "content_fixed",
    "content_auto",
    "chart_auto",
    "image_text_auto",
    "summary",
    "ending",
}
BUILTIN_TEMPLATE_ANNOTATION_TOKENS = [
    "标题",
    "标题1",
    "标题2",
    "编号",
    "小标题",
    "目录文字",
    "英文章节",
    "章节大标题",
    "章节副标题",
    "姓名",
    "职务",
    "单位",
    "日期",
    "职务、单位",
]


def _product_log_message(message: str) -> str:
    if message.startswith("Template ") and " loaded" in message:
        return "Template selected."
    return message


def _check_paper_ppt_agent_dependencies() -> None:
    diagnostics = get_academic_ppt_diagnostics()
    if diagnostics.missing_agent:
        raise RuntimeError("paper-ppt-agent local package was not found.")
    if diagnostics.missing_dependencies:
        raise RuntimeError(format_missing_dependency_message(diagnostics.missing_dependencies))


def _ensure_paper_ppt_agent_importable() -> Path:
    _check_paper_ppt_agent_dependencies()
    settings = get_settings()
    root = settings.paper_ppt_agent_root
    if not (root / "backend" / "orchestrator" / "pipeline.py").exists():
        raise RuntimeError("paper-ppt-agent local package was not found.")
    root_str = str(root)
    if root_str not in sys.path:
        sys.path.insert(0, root_str)
    return root


@contextmanager
def _paper_agent_import_context():
    """Import paper-ppt-agent without letting it read the NexusAI root .env."""
    previous_cwd = Path.cwd()
    engine_root = Path(__file__).resolve().parents[3]
    os.chdir(engine_root)
    try:
        yield
    finally:
        os.chdir(previous_cwd)


def _patch_paper_ppt_agent_runtime(task_dir: Path, model_bridge_url: str | None) -> None:
    _ensure_paper_ppt_agent_importable()

    with _paper_agent_import_context():
        import backend.llm.registry as registry
        import backend.orchestrator.svg_executor as svg_executor
        from backend.config import settings as paper_settings

    paper_root = get_settings().paper_ppt_agent_root
    runtime_dir = safe_task_child(task_dir, "checkpoints", "paper-runtime")
    workspaces_dir = safe_task_child(task_dir, "checkpoints", "paper-workspaces")

    paper_settings.assets_dir = paper_root / "assets"
    paper_settings.templates_dir = paper_root / "assets" / "templates"
    paper_settings.icons_dir = paper_root / "assets" / "icons"
    paper_settings.references_dir = paper_root / "assets" / "references"
    paper_settings.runtime_dir = runtime_dir
    paper_settings.workspaces_dir = workspaces_dir

    registry._PROVIDER_IMPORTS["nexus"] = ("app.core.model_bridge", "NexusModelBridgeProvider")
    registry._PROVIDER_INFO["nexus"] = NexusModelBridgeProvider(
        api_key="nexus-task",
        base_url=model_bridge_url or "",
    ).get_provider_info()

    original_create_provider = getattr(
        registry,
        "_nexusai_original_create_provider",
        registry.create_provider,
    )
    registry._nexusai_original_create_provider = original_create_provider

    def create_provider_with_nexus(
        name: str,
        api_key: str,
        *,
        base_url: str | None = None,
        deepseek_settings: dict | None = None,
        openai_settings: dict | None = None,
    ):
        if name == "nexus":
            return NexusModelBridgeProvider(api_key=api_key, base_url=base_url or model_bridge_url or "")
        return original_create_provider(
            name,
            api_key,
            base_url=base_url,
            deepseek_settings=deepseek_settings,
            openai_settings=openai_settings,
        )

    registry.create_provider = create_provider_with_nexus
    _patch_paper_ppt_agent_svg_executor(svg_executor)


def _patch_paper_ppt_agent_svg_executor(svg_executor: Any) -> None:
    if getattr(svg_executor, "_nexusai_chapter_detector_patched", False):
        return

    def detect_chapter_pages_from_intent(design_spec: str, total_pages: int) -> set[int]:
        chapter_pages: set[int] = set()
        slide_blocks = re.split(
            r"(?=^\s*#{2,6}\s+Slide\s+\d+\b|^\s*[-*]\s*Slide\s+\d+\b)",
            design_spec,
            flags=re.IGNORECASE | re.MULTILINE,
        )
        section_terms = (
            "section divider",
            "chapter divider",
            "section transition",
            "chapter transition",
            "part divider",
            "divider page",
            "\u7ae0\u8282\u8fc7\u6e21",
            "\u7ae0\u8282\u9875",
            "\u5206\u7ae0\u9875",
            "\u8fc7\u6e21\u9875",
        )
        body_terms = (
            "**content**",
            "content:",
            "visualization",
            "left-right split",
            "top-bottom split",
            "cards",
            "matrix",
            "process",
            "comparison",
            "\u6b63\u6587",
            "\u56fe\u8868",
            "\u65b9\u6cd5",
            "\u8bc1\u636e",
            "\u53d1\u73b0",
            "\u5bf9\u6bd4",
        )
        for block in slide_blocks:
            match = re.search(r"\bSlide\s+(\d+)\b", block, flags=re.IGNORECASE)
            if not match:
                continue
            page_num = int(match.group(1))
            if page_num <= 1 or page_num >= total_pages:
                continue
            head = "\n".join(block.splitlines()[:16])
            lower_head = head.lower()
            has_section_intent = any(term in lower_head or term in head for term in section_terms)
            has_body_intent = any(term in lower_head or term in head for term in body_terms)
            if has_section_intent and not has_body_intent:
                chapter_pages.add(page_num)
        return chapter_pages

    svg_executor._detect_chapter_pages = detect_chapter_pages_from_intent
    svg_executor._nexusai_chapter_detector_patched = True


def _escape_latex_text(text: str) -> str:
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    return "".join(replacements.get(ch, ch) for ch in text)


def _markdown_heading_title(text: str, fallback: str) -> str:
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        heading = re.match(r"^#{1,2}\s+(.+)$", stripped)
        if heading:
            title = heading.group(1).strip()
            if title:
                return title[:120]
        if len(stripped) >= 4 and not stripped.startswith(("-", "*", ">")):
            return stripped[:120]
    return fallback


def _text_to_tex(input_path: Path, task_dir: Path) -> Path:
    text = input_path.read_text(encoding="utf-8", errors="ignore")
    fallback_title = input_path.stem.replace("_", " ").strip() or "Academic Presentation"
    title = _markdown_heading_title(text, fallback_title)
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n+", text) if item.strip()]
    abstract = paragraphs[0][:1800] if paragraphs else text[:1800]
    body = paragraphs[1:] if len(paragraphs) > 1 else paragraphs
    section_chunks: list[str] = []

    for index, paragraph in enumerate(body[:24], start=1):
        heading = f"Section {index}"
        lines = paragraph.splitlines()
        first = lines[0].strip() if lines else ""
        if first.startswith("#"):
            heading = first.lstrip("#").strip()[:80] or heading
            paragraph = "\n".join(lines[1:]).strip() or paragraph
        section_chunks.append(
            f"\\section{{{_escape_latex_text(heading)}}}\n{_escape_latex_text(paragraph[:3000])}"
        )

    if not section_chunks:
        section_chunks.append(f"\\section{{Overview}}\n{_escape_latex_text(text[:3000])}")

    tex_path = safe_task_child(task_dir, "checkpoints", "normalized-source.tex")
    tex_path.write_text(
        "\n".join(
            [
                r"\documentclass{article}",
                f"\\title{{{_escape_latex_text(title)}}}",
                r"\begin{document}",
                r"\maketitle",
                r"\begin{abstract}",
                _escape_latex_text(abstract),
                r"\end{abstract}",
                *section_chunks,
                r"\end{document}",
            ]
        ),
        encoding="utf-8",
    )
    write_checkpoint(
        task_dir,
        "source-parsed",
        {"sourceType": "latex", "normalizedFrom": input_path.suffix.lower(), "fileName": tex_path.name},
    )
    return tex_path


def _pptx_to_tex(input_path: Path, task_dir: Path) -> Path:
    try:
        from pptx import Presentation

        presentation = Presentation(str(input_path))
        chunks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            if texts:
                chunks.append(f"# Slide {index}\n" + "\n".join(texts))
    except Exception as exc:
        raise RuntimeError(f"PPTX text extraction failed: {sanitize_message(exc)}") from exc

    extracted = "\n\n".join(chunks) or input_path.stem
    temp_md = safe_task_child(task_dir, "checkpoints", "normalized-pptx.md")
    temp_md.write_text(extracted, encoding="utf-8")
    return _text_to_tex(temp_md, task_dir)


def _normalize_input_for_paper_agent(input_path: Path, task_dir: Path) -> tuple[Path, str]:
    suffix = input_path.suffix.lower()
    if suffix == ".pdf":
        write_checkpoint(task_dir, "source-parsed", {"sourceType": "pdf", "fileName": input_path.name})
        return input_path, "pdf"
    if suffix in {".tex", ".zip", ".tgz"} or input_path.name.lower().endswith(".tar.gz"):
        write_checkpoint(task_dir, "source-parsed", {"sourceType": "latex", "fileName": input_path.name})
        return input_path, "latex"
    if suffix in {".txt", ".md", ".markdown"}:
        return _text_to_tex(input_path, task_dir), "latex"
    if suffix == ".pptx":
        return _pptx_to_tex(input_path, task_dir), "latex"
    raise RuntimeError(f"Unsupported academic-ppt input type: {suffix or input_path.name}")


def _canvas_format(settings: dict[str, Any]) -> str:
    return "ppt43" if settings.get("aspectRatio") == "4:3" else "ppt169"


def _detail_level(settings: dict[str, Any]) -> str:
    detail = settings.get("detailLevel")
    density = settings.get("informationDensity")
    if detail == "detailed" or density == "high":
        return "high"
    return "normal"


def _style(settings: dict[str, Any]) -> str:
    style = settings.get("templateStyle")
    if style == "blue_tech":
        return "tech"
    if style == "research_report":
        return "consulting"
    return "academic"


def _is_builtin_template_request(settings: dict[str, Any]) -> bool:
    return (
        str(settings.get("templateStyle") or "") == BUILTIN_TEMPLATE_ID
        or str(settings.get("templateId") or "") == BUILTIN_TEMPLATE_ID
    )


def _builtin_template_page_variants(metadata: dict[str, Any]) -> list[dict[str, Any]]:
    page_variants = metadata.get("pageVariants")
    if isinstance(page_variants, list):
        return [entry for entry in page_variants if isinstance(entry, dict)]
    variants = metadata.get("variants")
    if isinstance(variants, list):
        return [entry for entry in variants if isinstance(entry, dict)]
    return []


def _builtin_template_mode(metadata: dict[str, Any]) -> str:
    mode = str(metadata.get("mode") or "").strip().lower()
    if mode:
        return mode
    return "layout_blueprint" if metadata.get("source") == "builtin-pptx-template" else "theme_preset"


def _builtin_template_uses_layout_blueprint(settings: dict[str, Any]) -> bool:
    if not _is_builtin_template_request(settings):
        return False
    return _builtin_template_mode(_builtin_template_metadata()) != "theme_preset"


def _builtin_template_metadata() -> dict[str, Any]:
    if not BUILTIN_TEMPLATE_JSON.exists():
        raise RuntimeError(f"Built-in academic PPT template metadata is missing: {BUILTIN_TEMPLATE_ID}")
    try:
        metadata = json.loads(BUILTIN_TEMPLATE_JSON.read_text(encoding="utf-8"))
    except Exception as exc:
        raise RuntimeError(f"Built-in academic PPT template metadata is invalid: {BUILTIN_TEMPLATE_ID}") from exc
    if metadata.get("templateId") != BUILTIN_TEMPLATE_ID or metadata.get("sanitized") is not True:
        raise RuntimeError(f"Built-in academic PPT template is not sanitized: {BUILTIN_TEMPLATE_ID}")
    mode = _builtin_template_mode(metadata)
    source = str(metadata.get("source") or "")
    if mode == "theme_preset":
        if source not in {"theme-preset", "builtin-pptx-template"}:
            raise RuntimeError(f"Built-in academic PPT template source is invalid: {BUILTIN_TEMPLATE_ID}")
    elif source != "builtin-pptx-template":
        raise RuntimeError(f"Built-in academic PPT template source is invalid: {BUILTIN_TEMPLATE_ID}")
    font_policy = metadata.get("fontPolicy") or {}
    if font_policy.get("doNotExportFonts") is not True:
        raise RuntimeError(f"Built-in academic PPT template font policy is invalid: {BUILTIN_TEMPLATE_ID}")
    if metadata.get("templateFamily") not in {"cqupt-purple-academic", "cqet-purple-academic", "uestc-purple-theme-preset"}:
        raise RuntimeError(f"Built-in academic PPT template family is invalid: {BUILTIN_TEMPLATE_ID}")
    layout_policy = metadata.get("layoutPolicy") or {}
    required_layout_policy = ("disableGeneratedFooterPageNumber", "doNotMoveSchoolLogo", "doNotRedrawHeader")
    if not all(layout_policy.get(key) is True for key in required_layout_policy):
        raise RuntimeError(f"Built-in academic PPT template layout policy is invalid: {BUILTIN_TEMPLATE_ID}")
    school_policy = metadata.get("schoolIdentityPolicy") or {}
    required_school_policy = (
        "preserveExistingLogo",
        "preserveExistingSeal",
        "preserveExistingSchoolName",
        "doNotMoveLogo",
        "doNotResizeLogo",
        "doNotDuplicateLogo",
        "doNotGenerateNewLogo",
        "doNotUseLogoAsPlaceholder",
        "logoAreasAreForbiddenForContent",
    )
    if not all(school_policy.get(key) is True for key in required_school_policy):
        raise RuntimeError(f"Built-in academic PPT template school identity policy is invalid: {BUILTIN_TEMPLATE_ID}")
    theme = metadata.get("theme") or {}
    if not (theme.get("primary") or theme.get("primaryColor")):
        raise RuntimeError(f"Built-in academic PPT template theme is invalid: {BUILTIN_TEMPLATE_ID}")
    variants = _builtin_template_page_variants(metadata)
    placeholder_policy = metadata.get("placeholderPolicy") or {}
    if mode == "theme_preset":
        if (
            placeholder_policy.get("removePowerPointDefaultPrompts") is not True
            or placeholder_policy.get("removeUnfilledPlaceholders") is not True
        ):
            raise RuntimeError(f"Built-in academic PPT template placeholder policy is invalid: {BUILTIN_TEMPLATE_ID}")
        return metadata

    if not BUILTIN_TEMPLATE_PPTX.exists():
        raise RuntimeError(f"Built-in academic PPT template asset is missing: {BUILTIN_TEMPLATE_ID}")
    coordinate_system = metadata.get("coordinateSystem") or {}
    if (
        coordinate_system.get("type") != "inches"
        or float(coordinate_system.get("width") or 0) != 13.333
        or float(coordinate_system.get("height") or 0) != 7.5
    ):
        raise RuntimeError(f"Built-in academic PPT template coordinate system is invalid: {BUILTIN_TEMPLATE_ID}")
    if not isinstance(variants, list) or not variants:
        raise RuntimeError(f"Built-in academic PPT template blueprint variants are missing: {BUILTIN_TEMPLATE_ID}")
    roles = {str(entry.get("role") or "") for entry in variants if isinstance(entry, dict)}
    required_roles = {"cover", "toc", "section", "content_fixed", "content_auto", "chart_auto", "image_text_auto", "summary", "ending"}
    if not required_roles.issubset(roles):
        missing = ", ".join(sorted(required_roles - roles))
        raise RuntimeError(f"Built-in academic PPT template blueprint is missing role(s): {missing}")
    for entry in variants:
        if not isinstance(entry, dict):
            raise RuntimeError(f"Built-in academic PPT template blueprint entry is invalid: {BUILTIN_TEMPLATE_ID}")
        if not entry.get("variantId") or not entry.get("sourceFile") or not entry.get("sourceSlideIndex"):
            raise RuntimeError(f"Built-in academic PPT template variant identity is invalid: {BUILTIN_TEMPLATE_ID}")
        if not entry.get("templateSlideIndex") or entry.get("role") not in BUILTIN_TEMPLATE_ROLES:
            raise RuntimeError(f"Built-in academic PPT template variant role mapping is invalid: {BUILTIN_TEMPLATE_ID}")
    if (
        placeholder_policy.get("replaceOnlyKnownPlaceholders") is not True
        or placeholder_policy.get("removePowerPointDefaultPrompts") is not True
        or placeholder_policy.get("removeUnfilledPlaceholders") is not True
    ):
        raise RuntimeError(f"Built-in academic PPT template placeholder policy is invalid: {BUILTIN_TEMPLATE_ID}")
    return metadata


def _builtin_template_instruction(settings: dict[str, Any]) -> str:
    if not _is_builtin_template_request(settings):
        return ""
    metadata = _builtin_template_metadata()
    if _builtin_template_mode(metadata) == "theme_preset":
        theme = metadata.get("theme") or {}
        primary = theme.get("primary") or theme.get("primaryColor") or "#801C80"
        gradient_from = theme.get("gradientFrom") or theme.get("gradientStartColor") or "#811C81"
        gradient_to = theme.get("gradientTo") or theme.get("gradientEndColor") or "#9D229D"
        return "\n".join(
            [
                "## Built-in school academic theme preset",
                f"- Selected template_id: {BUILTIN_TEMPLATE_ID}. Treat it as a stable theme preset, not as a slide-by-slide template reconstruction task.",
                "- The final downloadable PPTX is the only authoritative output. Do not optimize for an SVG-only preview and do not describe a layout that the exported PPTX cannot match.",
                "- Use the native paper-ppt visual generation logic for layout, hierarchy, card composition, charts, and page-role structure. Do not flatten the deck into a fixed white-background template.",
                f"- Theme colors: primary {primary}; gradient {gradient_from} to {gradient_to}. Recolor the native paper-ppt theme into this Chongqing Polytechnic University of Electronic Technology purple palette.",
                "- Required page roles: cover, toc, section, content, chart/data, comparison, and summary. Do not generate an ending, thanks, back-cover, or closing-only slide for this template.",
                "- Cover page: may use the native paper-ppt cover composition, but it must not have a top decoration bar. Keep only the main title, one concise subtitle or metadata line, and the purple theme palette. Do not add template labels such as report cover.",
                "- TOC page: use the native agenda/list composition with clear numbering and generous spacing. Purple accents should replace the original palette.",
                "- Section page: use the native section/chapter composition with a strong visual anchor, concise section title, and the purple palette.",
                "- Content page: use native paper-ppt body layouts such as cards, columns, callouts, tables, timelines, and comparison structures. Do not turn every page into the same six-card layout.",
                "- Chart/data/comparison page: use native paper-ppt chart/table/matrix layouts, recolored to the purple palette with visible line weights and emphasized figures in the primary purple.",
                "- Summary/ending page: use native paper-ppt conclusion/ending structure, recolored to the purple palette and kept clean.",
                "- Keep each body slide to a scan-friendly amount of text: prefer 3-5 concise bullets or cards; split overly dense content instead of shrinking text below readable size.",
                "- Reserve the top-right school identity zone on non-cover slides. The adapter will place the verified template logo image there after generation.",
                "- Do not draw, synthesize, or write a school logo/school-name text mark yourself. Do not use text as a logo substitute.",
                "- Do not write 学术汇报, Academic Report, template labels, or any other text under or around the top-right logo.",
                "- Avoid large pale-purple circles, crossing-line decorations, decorative SVG motifs, and ornamental blobs. If decoration is needed, use only subtle lines, pale blocks, small dots, and card borders.",
                "- Do not use PowerPoint default placeholders, Click to add title/text prompts, {{...}} tokens, user annotation labels, or sample template text in the final PPTX.",
                "- Clean unfilled placeholders before export, keep the PPTX OpenXML package valid, and never embed fonts into the exported deck.",
            ]
        )
    theme = metadata.get("theme") or {}
    fallback_fonts = metadata.get("fontPolicy", {}).get("fallbackFonts") or ["SimSun", "SimHei", "Microsoft YaHei", "Arial"]
    variant_groups = _builtin_template_variant_groups(metadata)

    def mapped_slides(key: str) -> str:
        entries = variant_groups.get(key) or []
        values: list[str] = []
        for entry in entries:
            if not isinstance(entry, dict):
                continue
            values.append(
                f"{entry.get('variantId')}=template#{entry.get('templateSlideIndex')} "
                f"from {entry.get('sourceFile')}#{entry.get('sourceSlideIndex')}"
            )
        return "; ".join(values) or "mapped in template metadata"

    return "\n".join(
        [
            "## Built-in school academic template constraints",
            f"- Selected template_id: {BUILTIN_TEMPLATE_ID}. Use it as the visual reference for this deck.",
            "- The sanitized PPTX template keeps institutional masters, layouts, logo/background imagery, footer areas, page ratio, and embedded font references.",
            "- The current native pipeline consumes these constraints through the adapter style/theme layer; do not switch to a local TypeScript writer.",
            "- The final PPTX is recomposed by NexusAI from template blueprint variants; produce clean academic content and label slide role intent clearly.",
            "- Never copy source example text from the PPTX template. Generate all visible content from the uploaded academic material and replace placeholder roles with new content.",
            "- Do not invent author names, dates, or filename-derived paper titles. If the uploaded material has no author/date metadata, omit those fields rather than filling placeholders.",
            "- Use a formal school or institution report tone: clear title hierarchy, restrained decorations, strong whitespace, and content-first academic pages.",
            "- Preserve the original master header. Do not redraw, move, cover, or redesign the school logo, school name, top marker, or header area.",
            "- Preserve existing logo and school identity exactly. Never add a new logo, never move a logo, never duplicate a logo, and never use logo/seal/school-name areas as content slots.",
            "- Content must stay inside explicit body slots or each page's autoLayoutRegion, and it must avoid every forbiddenAreas rectangle in the template metadata.",
            "- Preserve original placeholder geometry: keep titles, subtitles, body copy, image regions, charts, and summary text inside the corresponding template regions.",
            "- Preserve original font hierarchy and approximate font sizes from the selected page type instead of inventing a new layout scale.",
            "- For generated Chinese text, use bold Songti-style headings where possible and Heiti-style body text; titles may wrap to two lines inside the mapped title slot.",
            "- Do not generate extra footer page numbers such as 01/05, 02/05, or slide-count counters. If the template has footer decoration, keep decoration only.",
            "- Use the dedicated TOC slide for the agenda/contents page; do not treat the TOC as a normal content slide.",
            "- Use section/chapter slides only for chapter transitions. Body pages must use content_fixed, content_auto, image_text_auto, or chart_auto layout logic.",
            "- Slide role rules: page 1 cover; page 2 TOC when there are at least 4 pages; section role only for chapter transition pages; body pages use content_fixed/content_auto/image_text_auto/chart_auto; the final page uses summary or ending.",
            "- The model must not choose template page numbers randomly. If a role is uncertain, mark it as content, not TOC/section/cover.",
            "- Clean every unfilled placeholder before export; no {{TITLE}}, {{BODY}}, {{KEY_POINTS}}, {{IMAGE}}, {{CHART}}, {{FOOTER}}, PowerPoint title prompt, or Click to add title may remain.",
            f"- Source template cover slides: {mapped_slides('cover')}; TOC slides: {mapped_slides('toc')}; section slides: {mapped_slides('section')}.",
            f"- Source template body slides: fixed={mapped_slides('content_fixed')}; auto={mapped_slides('content_auto')}; image/text={mapped_slides('image_text_auto')}; chart={mapped_slides('chart_auto')}; summary={mapped_slides('summary')}; ending={mapped_slides('ending')}.",
            f"- Preferred title color: {theme.get('titleColor') or '#3F3F3F'}; body color: {theme.get('bodyColor') or '#3F3F3F'}; accent color: {theme.get('accentColor') or '#801C80'}.",
            f"- Preferred palette: primary {theme.get('primaryColor') or '#801C80'}, gradient {theme.get('gradientStartColor') or '#9D229D'} to {theme.get('gradientEndColor') or '#801C80'}, background {theme.get('backgroundColor') or '#F2F2F2'}.",
            f"- Preserve the font feel by preferring these fallback font families when available: {', '.join(str(font) for font in fallback_fonts)}.",
            "- Use cover, agenda, section, content/chart, summary, and ending roles from the template metadata. Avoid turning most slides into section dividers.",
            "- If the built-in template asset or metadata cannot be verified, fail or mark the task degraded instead of claiming the template was applied.",
        ]
    )


def _template_id(settings: dict[str, Any]) -> str | None:
    if _is_builtin_template_request(settings):
        _builtin_template_metadata()
        return "academic_defense"
    mapping = {
        "academic_clean": "academic_defense",
        "blue_tech": "\u79d1\u6280\u84dd\u5546\u52a1",
        "research_report": "mckinsey",
        "course_presentation": "google_style",
    }
    return mapping.get(str(settings.get("templateStyle") or ""))


def _instruction(settings: dict[str, Any]) -> str:
    extras = str(settings.get("extraRequirements") or "").strip()
    flags = []
    if settings.get("enableDeepResearch"):
        flags.append("Use a deeper academic narrative with background, method, evidence, limitations, and conclusion.")
    if settings.get("enableExternalResearch"):
        flags.append("Use only server-provided context. Do not invent citations or external facts.")
    if settings.get("enableIconDecoration"):
        flags.append("Use restrained academic iconography where helpful.")
    if extras:
        flags.append(extras)
    flags.append(_slide_composition_guardrails(settings))
    flags.append(_visual_quality_guardrails(settings))
    builtin_instruction = _builtin_template_instruction(settings)
    if builtin_instruction:
        flags.append(builtin_instruction)
    return "\n".join(flags)


def _slide_composition_guardrails(settings: dict[str, Any]) -> str:
    target = max(1, int(settings.get("targetSlides") or 12))
    min_content = max(1, math.ceil(target * 0.55)) if target >= 8 else max(1, target - 3)
    max_sections = max(1, math.floor(target * 0.30)) if target >= 8 else 2
    max_outline = max(1, math.floor(target * 0.35)) if target >= 8 else 2
    return "\n".join(
        [
            "## NexusAI slide composition guardrails",
            f"- The requested deck has exactly {target} pages. Do not convert most pages into chapter dividers.",
            "- Recommended 12-page academic structure: 1 cover, 1 agenda, at most 2-3 section divider pages, 6-8 content/evidence pages, and 1 summary/ending page.",
            f"- For this request, produce at least {min_content} body content slides.",
            f"- Section divider/chapter-only slides must be no more than {max_sections}.",
            f"- Outline-only pages must be no more than {max_outline}.",
            "- A body content slide must include a clear title plus 3-5 concrete points from the manuscript.",
            "- Every body content slide must include at least one of: method, evidence, finding, figure, chart, comparison, implication, or conclusion.",
            "- Do not make a slide that only has a page number, PART label, chapter title, and subtitle unless it is explicitly one of the limited section divider pages.",
            "- For long papers, extract details from the manuscript into body slides instead of summarizing every chapter as a divider.",
            "- In the design_spec, label each page intent clearly as cover, agenda, section, content, evidence, method, findings, summary, or ending.",
        ]
    )


def _visual_quality_guardrails(settings: dict[str, Any]) -> str:
    language = "Chinese" if settings.get("outputLanguage") != "en" else "English"
    lines = [
        "## NexusAI visual quality guardrails",
        "These are design-quality constraints for the native paper-ppt-agent pipeline; keep all visible slide text in "
        f"{language}.",
        "- Readability comes first: every title, body line, caption, label, number, and footnote must remain readable on its background.",
        "- Use high-contrast title colors and medium-high-contrast body colors. Avoid dark text on deep blue or black backgrounds, pale gray text on light backgrounds, and low-opacity body text.",
        "- Cover pages need clear title, subtitle, and metadata hierarchy; the title should be the strongest text element.",
        "- Section-page large numbers are visual aids: keep them behind, beside, or clearly separated from the main title, never on top of readable text.",
        "- Agenda, summary, and closing pages should have balanced whitespace, consistent list spacing, and a clear visual anchor.",
        "- KPI cards, charts, flows, process diagrams, and comparison panels must be visibly separated from the page background with sufficient fill, stroke, or shadow contrast.",
        "- Keep information density balanced: avoid sparse empty slides, but do not pack text to the edges or exceed the content area.",
        "- Strict visual mode: do not emit a rule-fallback, plain white, or basic parsed-text design. If the visual strategy or SVG plan cannot be completed, fail cleanly so the task can be resumed.",
    ]
    if settings.get("templateStyle") == "blue_tech":
        lines.extend(
            [
                "## Dark blue template contrast rules",
                "- On dark or blue gradient backgrounds, title text must be #FFFFFF or a near-white high-contrast color.",
                "- On dark or blue gradient backgrounds, subtitle/body text must be #DCEBFF, #EAF4FF, or similarly light high-contrast text.",
                "- Do not place #0B1220, #172554, dark navy, or black text directly on dark blue backgrounds.",
                "- Agenda, summary, and section divider pages need separate color handling: agenda list text may use dark text only inside light cards; section and ending titles on dark backgrounds must stay white.",
                "- Large decorative numbers are decorative_number elements only: use low opacity and keep them behind title/subtitle text.",
                "- Section divider layer order is fixed: background, decorative_number, accent line or Part label, subtitle, title. The title is always the top readable layer.",
                "- Keep decorative numbers away from the title bounding box where possible; if they intersect, reduce opacity below 0.12 and draw title/subtitle after the number.",
                "- Part labels and accent rules may use cyan/blue accents, but title and subtitle on dark blue backgrounds must use the light text palette.",
                "- If a page uses a dark image or dark blue background, add a contrast panel or overlay before placing readable text.",
            ]
        )
    return "\n".join(lines)


def _blue_tech_style_overrides() -> dict[str, Any]:
    return {
        "title": "#FFFFFF",
        "subtitle": "#DCEBFF",
        "body": "#EAF4FF",
        "body_on_light": "#0B1220",
        "secondary_text": "#BFDBFE",
        "muted_text": "#BFD7FF",
        "accent": "#22D3EE",
        "accent_cyan": "#3DDCFF",
        "accent_alt": "#60A5FA",
        "accent_blue": "#5DB8FF",
        "card_fill": "#F8FAFC",
        "card_text": "#0B1220",
        "forbidden_dark_text_on_dark_background": ["#000000", "#020617", "#0B1220", "#0F172A", "#111827", "#172554", "#1E293B"],
        "decorative_number": {
            "fill": "#FFFFFF",
            "opacity": 0.08,
            "max_opacity_when_overlapping_text": 0.12,
            "layer": "behind_text",
            "z_index": 1,
        },
        "section_divider_layers": {
            "background": 0,
            "decorative_number": 1,
            "part_label": 2,
            "accent_rule": 2,
            "subtitle": 3,
            "title": 4,
        },
    }


def _visual_quality_style_overrides(settings: dict[str, Any]) -> dict[str, Any]:
    style = str(settings.get("templateStyle") or "academic_clean")
    density = "compact" if settings.get("informationDensity") == "high" else "normal"
    common = {
        "font_heading": "Aptos Display",
        "font_body": "Aptos",
        "cjk_heading": "SimSun",
        "cjk_body": "SimHei",
        "density": density,
        "readability_rules": [
            "Text must remain readable on its background.",
            "Use high contrast for titles and medium-high contrast for body text.",
            "Section-page large numbers must not overlap the main title.",
            "Cards, charts, and process panels need visible separation from the page background.",
            "For dark blue backgrounds, titles must be white and body text must use #DCEBFF or #EAF4FF.",
            "For section divider slides, draw decorative numbers before title/subtitle and keep the title visually on top.",
        ],
    }
    palettes = {
        "academic_clean": ["#F8FAFC", "#0F172A", "#1D4ED8", "#0F766E", "#334155", "#E2E8F0"],
        "blue_tech": ["#0B1220", "#F8FAFC", "#60A5FA", "#22D3EE", "#E2E8F0", "#172554"],
        "research_report": ["#FFFFFF", "#111827", "#2563EB", "#0F766E", "#374151", "#E5E7EB"],
        "course_presentation": ["#FFFFFF", "#111827", "#2563EB", "#F59E0B", "#374151", "#E5E7EB"],
        "school_academic_report": ["#FFFFFF", "#2F1B3A", "#801C80", "#811C81", "#9D229D", "#F3E6F5"],
    }
    result = {**common, "palette": palettes.get(style, palettes["academic_clean"])}
    if style == "blue_tech":
        result.update(
            {
                "dark_theme": True,
                "text_palette": _blue_tech_style_overrides(),
                "contrast_rules": {
                    "dark_background_title": "#FFFFFF",
                    "dark_background_body": "#DCEBFF",
                    "dark_background_body_alt": "#EAF4FF",
                    "dark_background_muted": "#BFD7FF",
                    "accent_cyan": "#3DDCFF",
                    "accent_blue": "#5DB8FF",
                    "forbidden_dark_text_on_dark_background": [
                        "#000000",
                        "#020617",
                        "#0B1220",
                        "#0F172A",
                        "#111827",
                        "#172554",
                        "#1E293B",
                    ],
                    "decorative_number": {
                        "fill": "#FFFFFF",
                        "opacity": 0.08,
                        "max_opacity_when_overlapping_text": 0.12,
                        "layer": "behind_text",
                        "z_index": 1,
                    },
                    "section_divider_layers": {
                        "background": 0,
                        "decorative_number": 1,
                        "part_label": 2,
                        "accent_rule": 2,
                        "subtitle": 3,
                        "title": 4,
                    },
                    "agenda_card_text": "#0B1220",
                    "summary_card_text": "#0B1220",
                },
            }
        )
    if style == BUILTIN_TEMPLATE_ID:
        metadata = _builtin_template_metadata()
        theme = metadata.get("theme") or {}
        if _builtin_template_mode(metadata) == "theme_preset":
            result.update(
                {
                    "builtin_template_id": BUILTIN_TEMPLATE_ID,
                    "builtin_template_asset_verified": True,
                    "template_source": "theme_preset",
                    "font_heading": "SimSun",
                    "font_body": "SimHei",
                    "cjk_heading": "SimSun",
                    "cjk_body": "SimHei",
                    "theme": {
                        "mode": "school_theme_preset",
                        "primary": theme.get("primary") or theme.get("primaryColor") or "#801C80",
                        "secondary": theme.get("secondary") or theme.get("secondaryColor") or "#9D229D",
                        "gradient_from": theme.get("gradientFrom") or theme.get("gradientStartColor") or "#811C81",
                        "gradient_to": theme.get("gradientTo") or theme.get("gradientEndColor") or "#9D229D",
                        "background": theme.get("background") or theme.get("backgroundColor") or "#FFFFFF",
                        "surface": theme.get("surfaceColor") or "#FFFFFF",
                        "title": theme.get("titleColor") or "#2F1B3A",
                        "body": theme.get("bodyColor") or "#3F3F3F",
                        "accent": theme.get("accentColor") or "#801C80",
                    },
                    "role_layout_guidance": {
                        "cover": "use the native paper-ppt cover composition with no top decoration bar, minimal text, polished title hierarchy, and the purple palette",
                        "toc": "use the native paper-ppt agenda/list composition with clear numbering and generous spacing",
                        "section": "use the native paper-ppt section composition with a strong visual anchor and concise title",
                        "content": "use native paper-ppt content layouts: cards, columns, callouts, tables, timelines, and comparison structures",
                        "chart": "use native paper-ppt chart/table/matrix layouts with the purple palette and visible line weights",
                        "comparison": "use native paper-ppt balanced comparison layouts with consistent card heights and clear labels",
                        "summary": "use native paper-ppt conclusion structure with clean evidence and next-step hierarchy",
                        "ending": "use native paper-ppt closing structure with restrained purple accents",
                    },
                    "template_family": metadata.get("templateFamily"),
                    "coordinate_system": metadata.get("coordinateSystem") or {},
                    "layout_types": metadata.get("layoutTypes") or {},
                    "school_identity_policy": metadata.get("schoolIdentityPolicy") or {},
                    "layout_policy": metadata.get("layoutPolicy") or {},
                    "placeholder_policy": metadata.get("placeholderPolicy") or {},
                    "logo_policy": {
                        "position": "top-right",
                        "source": BUILTIN_TEMPLATE_LOGO_MEDIA,
                        "preserve_existing_identity": True,
                        "do_not_generate_logo": True,
                        "do_not_move_logo": True,
                        "do_not_duplicate_logo": True,
                        "do_not_write_text_fallback": True,
                        "reserved_zone": {"x": 10.0, "y": 0.0, "w": 3.05, "h": 0.9},
                    },
                    "slide_role_presets": [
                        "cover",
                        "toc",
                        "section",
                        "content",
                        "chart",
                        "comparison",
                        "summary",
                        "ending",
                    ],
                    "page_number_policy": "disabled",
                    "preserve_embedded_fonts": False,
                    "do_not_export_fonts": True,
                    "readability_rules": [
                        *common["readability_rules"],
                        "Keep the native paper-ppt layout logic and only recolor the generated design into the school purple palette.",
                        "Do not force a uniform top-bar/card template across all slides.",
                        "Reserve the top-right school identity zone and keep all generated content outside that forbidden area.",
                        "Never invent, redraw, duplicate, or write logos, seals, school-name text marks, or 学术汇报 labels; the adapter places the verified logo image.",
                        "Do not create a closing-only, thanks, back-cover, or ending slide; use the final slide as a normal summary/content slide.",
                        "Prefer compact cards, structured comparisons, tables, and summary callouts over free-form text walls.",
                        "Use 3-5 concise bullets or cards per slide where possible; split dense material instead of overcrowding.",
                        "Keep chart and comparison strokes visible in Microsoft PowerPoint; avoid hairline-only separators.",
                        "Do not add generated footer page numbers.",
                        "Do not include sanitized placeholders, PowerPoint default prompts, or template sample text in the final deck.",
                    ],
                }
            )
        else:
            result.update(
                {
                    "builtin_template_id": BUILTIN_TEMPLATE_ID,
                    "builtin_template_asset_verified": True,
                    "template_source": "builtin-pptx-template",
                    "font_heading": "SimSun",
                    "font_body": "SimHei",
                    "cjk_heading": "SimSun",
                    "cjk_body": "SimHei",
                    "theme": {
                        "primary": theme.get("primaryColor") or "#801C80",
                        "secondary": theme.get("secondaryColor") or "#9D229D",
                        "background": theme.get("backgroundColor") or "#F2F2F2",
                        "title": theme.get("titleColor") or "#3F3F3F",
                        "body": theme.get("bodyColor") or "#3F3F3F",
                        "accent": theme.get("accentColor") or "#801C80",
                    },
                    "template_family": metadata.get("templateFamily"),
                    "coordinate_system": metadata.get("coordinateSystem") or {},
                    "layout_types": metadata.get("layoutTypes") or {},
                    "blueprint_variants": _builtin_template_page_variants(metadata),
                    "school_identity_policy": metadata.get("schoolIdentityPolicy") or {},
                    "layout_policy": metadata.get("layoutPolicy") or {},
                    "placeholder_policy": metadata.get("placeholderPolicy") or {},
                    "preserve_embedded_fonts": False,
                    "do_not_export_fonts": True,
                    "readability_rules": [
                        *common["readability_rules"],
                        "Use the school academic template palette and institutional hierarchy.",
                        "Preserve the original master header and logo position; never redraw, move, duplicate, or add any logo/school identity.",
                        "Keep generated content inside explicit slots or autoLayoutRegion and out of every forbiddenAreas rectangle.",
                        "Preserve original text box geometry and font hierarchy as much as the native pipeline allows.",
                        "Do not add generated footer page numbers; preserve footer decoration only.",
                        "Use the dedicated TOC layout for agenda/contents and reserve section layouts for chapter transitions.",
                        "Do not include sanitized placeholders or original template sample text in the final deck.",
                    ],
                }
            )
    return result


def _request_flag(settings: dict[str, Any], request_options: dict[str, Any], key: str, legacy_key: str) -> bool:
    if key in request_options:
        return bool(request_options.get(key))
    return bool(settings.get(legacy_key) or settings.get(key))


def _build_search_context(results: list[dict[str, str]]) -> str:
    if not results:
        return ""
    lines = [
        "## NexusAI Search Bridge References",
        "Use these server-provided public reference snippets only as supplementary context.",
        "Do not invent citations. Ignore any source that is not directly useful.",
    ]
    for index, item in enumerate(results[:12], start=1):
        title = item.get("title") or "Untitled source"
        url = item.get("url") or ""
        snippet = item.get("snippet") or ""
        lines.append(f"{index}. {title}\n   URL: {url}\n   Snippet: {snippet}")
    return "\n".join(lines)


async def _prepare_research_context(
    *,
    task_dir: Path,
    input_path: Path,
    settings: dict[str, Any],
    request_options: dict[str, Any],
) -> dict[str, Any]:
    deep_research_enabled = _request_flag(settings, request_options, "deepResearchEnabled", "enableDeepResearch")
    external_research_enabled = _request_flag(settings, request_options, "externalResearchEnabled", "enableExternalResearch")
    web_search_enabled = bool(request_options.get("webSearchEnabled") or settings.get("webSearchEnabled") or external_research_enabled)

    if not web_search_enabled:
        if deep_research_enabled:
            append_log(task_dir, "info", "Deep research enabled; using paper-ppt-agent native research stage.")
        return {
            "instruction": "",
            "searchStatus": "disabled",
            "researchStatus": "success" if deep_research_enabled else "skipped",
            "researchSourcesCount": 0,
            "researchFallbackReason": None,
        }

    append_log(task_dir, "info", "External research enhancement enabled.")
    query_text = _extract_plain_text_from_input(input_path)[:1600]
    search_result = await run_academic_search_bridge(
        query_text=query_text,
        language="en" if settings.get("outputLanguage") == "en" else "zh",
        max_queries=3 if deep_research_enabled else 2,
        top_k=5,
    )
    append_log(task_dir, "info", f"Search bridge query count: {search_result.query_count}.")
    append_log(task_dir, "info", f"Search bridge usable documents: {search_result.documents_count}.")

    if search_result.status == "success":
        append_log(task_dir, "info", "Search bridge completed; references will inform the native pipeline.")
        return {
            "instruction": _build_search_context(search_result.results),
            "searchStatus": "success",
            "researchStatus": "success",
            "researchSourcesCount": search_result.documents_count,
            "researchFallbackReason": None,
        }

    fallback_reason = search_result.error_summary or "Search bridge returned no usable public sources."
    append_log(task_dir, "warn", "Search bridge degraded; continuing without external references.")
    return {
        "instruction": "",
        "searchStatus": "degraded",
        "researchStatus": "degraded",
        "researchSourcesCount": search_result.documents_count,
        "researchFallbackReason": fallback_reason,
    }


def _safe_output_file(task_dir: Path, raw_output_path: str | None) -> tuple[Path, int]:
    output_path = safe_task_child(task_dir, "outputs", "academic-ppt-result.pptx")
    if not raw_output_path:
        raise RuntimeError("paper-ppt-agent did not return a PPTX output.")

    source = Path(raw_output_path).resolve()
    paper_workspace = safe_task_child(task_dir, "checkpoints", "paper-workspaces")
    if not source.exists() or not source.is_file():
        raise RuntimeError("paper-ppt-agent PPTX output is missing.")
    try:
        source.relative_to(paper_workspace.resolve())
    except ValueError as exc:
        raise RuntimeError("paper-ppt-agent output path escaped task workspace.") from exc

    shutil.copy2(source, output_path)
    return output_path, output_path.stat().st_size


def _builtin_generated_text_replacements(settings: dict[str, Any]) -> dict[str, str]:
    if not _is_builtin_template_request(settings):
        return {}
    replacements = {
        "作者：NexusAI": "",
        "NexusAI": "",
        "论文标题：": "主题：",
        "论文标题": "主题",
        "2025-08-08": "",
        "学校简介": "",
        "学校概况": "",
        "重电实践": "",
        "申报准备": "",
        "建设探索": "",
        "时代背景": "",
        "重电举措": "",
        "未来计划": "",
        "未来规划": "",
        "国家战略有要求": "",
        "数字化重塑": "",
        "新双高、新内涵": "",
        "双高计划": "",
        "DeepSeek": "",
        "OBE+AI": "",
        "达特茅斯": "",
        "Artificial Intelligence": "",
        "单击此处添加标题": "",
        "单击添加标题": "",
        "单击此处添加文本": "",
        "Click to add title": "",
        "Click to add text": "",
    }
    replacements.update({token: "" for token in BUILTIN_TEMPLATE_ANNOTATION_TOKENS})
    placeholders = _builtin_placeholder_tokens()
    replacements.update({token: "" for token in placeholders})
    replacements.update({prompt: "" for prompt in _builtin_powerpoint_default_prompts()})
    return replacements


def _builtin_placeholder_tokens() -> list[str]:
    try:
        metadata = _builtin_template_metadata()
    except Exception:
        return ["{{TITLE}}", "{{BODY}}", "{{KEY_POINTS}}", "{{IMAGE}}", "{{CHART}}", "{{FOOTER}}"]
    policy = metadata.get("placeholderPolicy") or {}
    tokens = policy.get("knownPlaceholders")
    if isinstance(tokens, list) and tokens:
        return [str(token) for token in tokens if str(token).strip()]
    return ["{{TITLE}}", "{{BODY}}", "{{KEY_POINTS}}", "{{IMAGE}}", "{{CHART}}", "{{FOOTER}}"]


def _builtin_powerpoint_default_prompts() -> list[str]:
    prompt_prefix = "".join(chr(value) for value in (0x5355, 0x51FB, 0x6B64, 0x5904, 0x6DFB, 0x52A0))
    prompts = [
        prompt_prefix,
        "单击此处添加标题",
        "单击此处添加项正文",
        "单击此处添加正文",
        "Click to add title",
        "Click to add text",
    ]
    try:
        metadata = _builtin_template_metadata()
    except Exception:
        return prompts
    policy = metadata.get("placeholderPolicy") or {}
    configured = policy.get("defaultPrompts")
    if isinstance(configured, list):
        prompts.extend(str(prompt) for prompt in configured if str(prompt).strip())
    return list(dict.fromkeys(prompts))


def _sanitize_builtin_template_svg_text(project_dir: str | None, settings: dict[str, Any]) -> int:
    replacements = _builtin_generated_text_replacements(settings)
    if not replacements or not project_dir:
        return 0
    changed = 0
    for svg_file in _svg_final_files(project_dir):
        svg = svg_file.read_text(encoding="utf-8", errors="ignore")
        next_svg = svg
        for old, new in replacements.items():
            next_svg = next_svg.replace(old, new)
        if next_svg != svg:
            svg_file.write_text(next_svg, encoding="utf-8")
            changed += 1
    return changed


def _replace_pptx_text_literals(pptx_path: Path, replacements: dict[str, str]) -> int:
    if not replacements:
        return 0
    temp_path = pptx_path.with_suffix(f"{pptx_path.suffix}.cleaned.tmp")
    changed = 0
    with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            payload = source.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/") or item.filename.startswith("docProps/")
            ):
                text = payload.decode("utf-8", errors="ignore")
                original = text
                for old, new in replacements.items():
                    text = text.replace(old, new)
                if text != original:
                    changed += 1
                    payload = text.encode("utf-8")
            target.writestr(item, payload)
    temp_path.replace(pptx_path)
    return changed


def _clean_builtin_template_shape_text(value: str) -> bool:
    text = (value or "").replace("\x0b", "").strip()
    if not text:
        return False
    if text in BUILTIN_TEMPLATE_ANNOTATION_TOKENS:
        return True
    if "{{" in text or "}}" in text:
        return True
    if any(prompt in text for prompt in _builtin_powerpoint_default_prompts()):
        return True
    return bool(text) and set(text) <= {"|"}


def _set_pptx_shape_text(
    shape: Any,
    text: str,
    *,
    font_size: int | None = None,
    font_family: str | None = None,
    color: str = "3F3F3F",
    bold: bool = False,
    alignment: Any | None = None,
) -> None:
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    if not lines:
        lines = [""]
    paragraph = text_frame.paragraphs[0]
    paragraph.text = lines[0]
    for line in lines[1:]:
        next_paragraph = text_frame.add_paragraph()
        next_paragraph.text = line
    try:
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.util import Pt

        rgb = RGBColor.from_string(color)
        resolved_font = str(font_family or "SimHei").strip() or "SimHei"

        def apply_typeface(run: Any) -> None:
            run.font.name = resolved_font
            r_pr = run._r.get_or_add_rPr()
            for tag_name in ("a:latin", "a:ea", "a:cs"):
                element = r_pr.find(qn(tag_name))
                if element is None:
                    element = OxmlElement(tag_name)
                    r_pr.append(element)
                element.set("typeface", resolved_font)

        for paragraph in text_frame.paragraphs:
            if alignment is not None:
                paragraph.alignment = alignment
            for run in paragraph.runs:
                apply_typeface(run)
                if font_size:
                    run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = rgb
    except Exception:
        return


def _add_builtin_template_textbox(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 16,
    font_family: str | None = None,
    color: str = "3F3F3F",
    bold: bool = False,
    alignment: Any | None = None,
) -> None:
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _set_pptx_shape_text(
        shape,
        text,
        font_size=font_size,
        font_family=font_family,
        color=color,
        bold=bold,
        alignment=alignment,
    )


def _theme_color(metadata: dict[str, Any], key: str, fallback: str) -> str:
    theme = metadata.get("theme") or {}
    aliases = {
        "primary": ("primary", "primaryColor"),
        "gradient_from": ("gradientFrom", "gradientStartColor"),
        "gradient_to": ("gradientTo", "gradientEndColor"),
        "background": ("background", "backgroundColor"),
        "title": ("titleColor", "title"),
        "body": ("bodyColor", "body"),
    }
    for name in aliases.get(key, (key,)):
        value = str(theme.get(name) or "").strip()
        if value:
            return value.lstrip("#")
    return fallback.lstrip("#")


def _theme_preset_shape_fill(shape: Any, color: str, transparency: float | None = None) -> None:
    try:
        from pptx.dml.color import RGBColor

        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
        if transparency is not None:
            shape.fill.transparency = max(0.0, min(1.0, float(transparency)))
        shape.line.fill.background()
    except Exception:
        return


def _theme_preset_add_rect(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    color: str,
    *,
    transparency: float | None = None,
) -> Any:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    _theme_preset_shape_fill(shape, color, transparency=transparency)
    return shape


def _theme_preset_add_gradient_bar(
    slide: Any,
    *,
    y: float = 0.0,
    height: float = 0.72,
    start_color: str = "811C81",
    end_color: str = "9D229D",
    slide_width: float = 13.333,
) -> Any:
    from pptx.util import Inches

    start_hex = start_color.lstrip("#").upper()
    end_hex = end_color.lstrip("#").upper()
    cache_key = hashlib.sha1(f"{start_hex}-{end_hex}-{height}".encode("ascii", "ignore")).hexdigest()[:12]
    gradient_path = Path(os.getenv("TEMP") or os.getenv("TMP") or ".") / f"academic-ppt-gradient-{cache_key}.png"
    if not gradient_path.exists():
        try:
            from PIL import Image

            width_px = 1600
            height_px = max(8, int(round(width_px * height / max(slide_width, 0.1))))
            start_rgb = tuple(int(start_hex[index : index + 2], 16) for index in (0, 2, 4))
            end_rgb = tuple(int(end_hex[index : index + 2], 16) for index in (0, 2, 4))
            image = Image.new("RGB", (width_px, height_px))
            pixels = image.load()
            for x_pos in range(width_px):
                ratio = x_pos / max(width_px - 1, 1)
                color = tuple(int(start_rgb[i] + (end_rgb[i] - start_rgb[i]) * ratio) for i in range(3))
                for y_pos in range(height_px):
                    pixels[x_pos, y_pos] = color
            image.save(gradient_path, format="PNG")
        except Exception:
            return _theme_preset_add_rect(slide, 0, y, slide_width, height, start_hex)
    picture = slide.shapes.add_picture(
        str(gradient_path),
        Inches(0),
        Inches(y),
        width=Inches(slide_width),
        height=Inches(height),
    )
    picture.name = "cqet_header_gradient"
    return picture


def _theme_preset_title_from_slide(slide: Any, fallback: str) -> str:
    candidates: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = _normalized_text_line(getattr(shape, "text", ""))
        if not text:
            continue
        if "{{" in text or "}}" in text:
            continue
        if any(prompt in text for prompt in _builtin_powerpoint_default_prompts()):
            continue
        if text in {"学术报告", "学术汇报", "目录", "CONTENTS", "SECTION"}:
            continue
        if text.startswith("数据来源") or text.lower().startswith("source:"):
            continue
        candidates.append(re.sub(r"^\d{1,2}\s*", "", text).strip())
    return (candidates[0] if candidates else fallback)[:64]


def _theme_preset_slide_texts(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = _normalized_text_line(getattr(shape, "text", ""))
        if text:
            texts.append(text)
    return _clean_generated_text_items(texts)


def _theme_preset_clear_text_shapes(slide: Any) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        _set_pptx_shape_text(shape, "")


def _theme_preset_add_header(
    slide: Any,
    title: str,
    metadata: dict[str, Any],
    *,
    slide_width: float = 13.333,
) -> None:
    from pptx.enum.text import PP_ALIGN

    gradient_from = _theme_color(metadata, "gradient_from", "811C81")
    gradient_to = _theme_color(metadata, "gradient_to", "9D229D")
    _theme_preset_add_gradient_bar(slide, height=0.72, start_color=gradient_from, end_color=gradient_to, slide_width=slide_width)
    _add_builtin_template_textbox(
        slide,
        0.58,
        0.16,
        8.4,
        0.36,
        title,
        font_size=17,
        font_family="SimHei",
        color="FFFFFF",
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def _theme_preset_apply_cover(slide: Any, metadata: dict[str, Any]) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    texts = _theme_preset_slide_texts(slide)
    title = (texts[0] if texts else "学术汇报")[:72]
    subtitle = (texts[1] if len(texts) > 1 else "")[:52]
    title_color = _theme_color(metadata, "title", "2F1B3A")
    title_weight = sum(1.0 if ord(char) > 127 else 0.58 for char in title)
    if title_weight <= 18:
        title_font_size = 48
        title_y = 2.78
        title_height = 0.92
        subtitle_y = 4.12
    elif title_weight <= 34:
        title_font_size = 44
        title_y = 2.66
        title_height = 1.18
        subtitle_y = 4.23
    else:
        title_font_size = 40
        title_y = 2.54
        title_height = 1.42
        subtitle_y = 4.36

    for shape in list(slide.shapes):
        _remove_shape_from_slide(shape)
    if BUILTIN_TEMPLATE_COVER_BACKGROUND.exists():
        background = slide.shapes.add_picture(
            str(BUILTIN_TEMPLATE_COVER_BACKGROUND),
            Inches(0),
            Inches(0),
            width=Inches(13.333),
            height=Inches(7.5),
        )
        background.name = "cqet_cover_background"
    else:
        _theme_preset_add_rect(slide, 0, 0, 13.333, 7.5, "FFFFFF")
    _add_builtin_template_textbox(
        slide,
        0.8,
        title_y,
        11.733,
        title_height,
        title,
        font_size=title_font_size,
        font_family="SimHei",
        color=title_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if subtitle:
        _add_builtin_template_textbox(
            slide,
            0.8,
            subtitle_y,
            11.733,
            0.35,
            subtitle,
            font_size=15,
            font_family="SimHei",
            color="6B556F",
            alignment=PP_ALIGN.CENTER,
        )


def _theme_preset_add_toc_body(slide: Any, texts: list[str], metadata: dict[str, Any]) -> None:
    from pptx.enum.text import PP_ALIGN

    primary = _theme_color(metadata, "primary", "801C80")
    items = [item for item in texts[1:] if item and item not in {"目录", "CONTENTS"}][:5]
    if not items:
        items = ["研究背景与评价口径", "核心发现与关键证据", "变化趋势与专业对比", "结论建议与后续行动"]
    for index, item in enumerate(items, start=1):
        y = 1.42 + (index - 1) * 0.86
        _theme_preset_add_rect(slide, 0.9, y, 11.45, 0.58, "FFFFFF")
        _theme_preset_add_rect(slide, 0.9, y, 0.08, 0.58, primary)
        _theme_preset_add_rect(slide, 1.18, y + 0.14, 0.38, 0.3, primary)
        _add_builtin_template_textbox(
            slide,
            1.18,
            y + 0.13,
            0.38,
            0.22,
            f"{index:02d}",
            font_size=9,
            font_family="SimHei",
            color="FFFFFF",
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        _add_builtin_template_textbox(
            slide,
            1.82,
            y + 0.12,
            9.7,
            0.3,
            item[:52],
            font_size=15,
            font_family="SimHei",
            color="2F1B3A",
            bold=True,
        )


def _theme_preset_add_default_body(slide: Any, texts: list[str], metadata: dict[str, Any], role: str) -> None:
    from pptx.enum.text import PP_ALIGN

    primary = _theme_color(metadata, "primary", "801C80")
    body = [item for item in texts[1:] if item][:6]
    if not body:
        body = ["围绕研究问题提炼关键发现", "结合证据材料形成结构化判断", "服务后续专业建设与决策优化"]
    if role == "section":
        _theme_preset_add_rect(slide, 0.75, 1.55, 3.15, 4.65, primary, transparency=0.06)
        _add_builtin_template_textbox(slide, 1.04, 2.18, 2.45, 0.42, "SECTION", font_size=18, font_family="SimHei", color=primary, bold=True)
        _add_builtin_template_textbox(slide, 4.35, 2.55, 7.65, 0.85, texts[0][:48] if texts else "章节概览", font_size=28, font_family="SimHei", color="2F1B3A", bold=True)
        _add_builtin_template_textbox(slide, 4.38, 3.58, 6.8, 0.42, body[0][:80], font_size=15, font_family="SimHei", color="6B556F")
        return
    if role == "summary":
        _theme_preset_add_rect(slide, 0.95, 1.48, 11.45, 1.05, "F6ECF8")
        _add_builtin_template_textbox(slide, 1.25, 1.75, 10.8, 0.34, body[0][:86], font_size=17, font_family="SimHei", color="2F1B3A", bold=True, alignment=PP_ALIGN.CENTER)
        items = body[1:5] or body[:4]
        for index, item in enumerate(items, start=1):
            y = 3.0 + (index - 1) * 0.62
            _theme_preset_add_rect(slide, 1.35, y + 0.07, 0.16, 0.16, primary)
            _add_builtin_template_textbox(slide, 1.75, y - 0.02, 9.8, 0.32, item[:100], font_size=13, font_family="SimHei", color="3F3F3F")
        return
    columns = 2 if len(body) <= 4 else 3
    card_w = 5.55 if columns == 2 else 3.62
    start_x = 0.85
    gap = 0.35
    for index, item in enumerate(body[:6]):
        col = index % columns
        row = index // columns
        x = start_x + col * (card_w + gap)
        y = 1.45 + row * 1.72
        _theme_preset_add_rect(slide, x, y, card_w, 1.28, "FFFFFF")
        _theme_preset_add_rect(slide, x, y, 0.08, 1.28, primary)
        _add_builtin_template_textbox(slide, x + 0.28, y + 0.2, card_w - 0.55, 0.26, item[:34], font_size=14, font_family="SimHei", color=primary, bold=True)
        if len(item) > 34:
            _add_builtin_template_textbox(slide, x + 0.28, y + 0.58, card_w - 0.55, 0.42, item[34:112], font_size=11, font_family="SimHei", color="4B5563")


def _remove_shape_from_slide(shape: Any) -> bool:
    try:
        element = shape._element
        parent = element.getparent()
        if parent is None:
            return False
        parent.remove(element)
        return True
    except Exception:
        return False


def _send_shape_to_back(shape: Any) -> bool:
    try:
        element = shape._element
        parent = element.getparent()
        if parent is None:
            return False
        parent.remove(element)
        parent.insert(0, element)
        return True
    except Exception:
        return False


def _bring_shape_to_front(shape: Any) -> bool:
    try:
        element = shape._element
        parent = element.getparent()
        if parent is None:
            return False
        parent.remove(element)
        parent.append(element)
        return True
    except Exception:
        return False


def _shape_inches(value: Any) -> float:
    try:
        return float(value) / 914400.0
    except Exception:
        return 0.0


def _shape_text_value(shape: Any) -> str:
    return _normalized_text_line(getattr(shape, "text", "")) if hasattr(shape, "text") else ""


def _shape_fill_rgb(shape: Any) -> str | None:
    try:
        rgb = shape.fill.fore_color.rgb
        return str(rgb).upper() if rgb is not None else None
    except Exception:
        return None


def _shape_bounds(shape: Any) -> tuple[float, float, float, float]:
    left = _shape_inches(getattr(shape, "left", 0))
    top = _shape_inches(getattr(shape, "top", 0))
    width = _shape_inches(getattr(shape, "width", 0))
    height = _shape_inches(getattr(shape, "height", 0))
    return left, top, width, height


def _shape_intersects(x: float, y: float, width: float, height: float, area: tuple[float, float, float, float]) -> bool:
    ax, ay, aw, ah = area
    return x < ax + aw and x + width > ax and y < ay + ah and y + height > ay


def _remove_theme_preset_logo_layer(slide: Any) -> int:
    removed = 0
    school_mark_texts = {
        "重庆电子科技职业大学",
        "CHONGQING POLYTECHNIC UNIVERSITY OF ELECTRONIC TECHNOLOGY",
        "学术汇报",
        "汇报分析",
        "ACADEMIC REPORT",
        "REPORT ANALYSIS",
    }
    logo_area = (9.75, 0.0, 3.45, 0.95)
    for shape in list(slide.shapes):
        shape_name = str(getattr(shape, "name", "") or "")
        text = _shape_text_value(shape)
        x, y, width, height = _shape_bounds(shape)
        in_top_right = _shape_intersects(x, y, width, height, logo_area)
        contains_school_mark = any(mark and mark in text for mark in school_mark_texts)
        if shape_name == "cqet_image9_logo" or (in_top_right and (text or contains_school_mark or width <= 3.45)):
            if _remove_shape_from_slide(shape):
                removed += 1
    return removed


def _remove_theme_preset_cover_header(slide: Any) -> int:
    removed = 0
    header_area = (0.0, 0.0, 13.333, 1.1)
    for shape in list(slide.shapes):
        x, y, width, height = _shape_bounds(shape)
        if not _shape_intersects(x, y, width, height, header_area):
            continue
        fill_rgb = _shape_fill_rgb(shape)
        text = _shape_text_value(shape)
        if (
            str(getattr(shape, "name", "") or "") == "cqet_image9_logo"
            or fill_rgb in {"2F1B3A", "2F1838", "321B3D", "3A1D43", "801C80", "811C81", "9D229D"}
            or "重庆电子科技职业大学" in text
            or "CHONGQING POLYTECHNIC" in text.upper()
            or text in {"学术汇报", "ACADEMIC REPORT"}
        ):
            if _remove_shape_from_slide(shape):
                removed += 1
    return removed


def _remove_theme_preset_cover_excess_text(slide: Any) -> int:
    text_shapes: list[tuple[Any, str, float, float, float, float]] = []
    for shape in slide.shapes:
        text = _shape_text_value(shape)
        if not text:
            continue
        x, y, width, height = _shape_bounds(shape)
        if y < 1.0 or y > 6.8:
            continue
        if "来源" in text or "第" in text and "页" in text:
            continue
        text_shapes.append((shape, text, x, y, width, height))
    protected = set()
    for item in sorted(text_shapes, key=lambda entry: len(entry[1]), reverse=True)[:2]:
        if item[3] >= 4.1:
            continue
        protected.add(id(item[0]))
    removed = 0
    for shape, text, x, y, width, height in text_shapes:
        if id(shape) in protected:
            continue
        if y >= 4.1 or len(text) <= 48:
            if _remove_shape_from_slide(shape):
                removed += 1
    for shape in list(slide.shapes):
        if id(shape) in protected:
            continue
        x, y, width, height = _shape_bounds(shape)
        if y >= 4.65 and y <= 6.95:
            if _remove_shape_from_slide(shape):
                removed += 1
    return removed


def _remove_theme_preset_top_bar(slide: Any) -> int:
    removed = 0
    header_area = (0.0, 0.0, 13.333, 0.85)
    for shape in list(slide.shapes):
        if str(getattr(shape, "name", "") or "") == "cqet_image9_logo":
            continue
        x, y, width, height = _shape_bounds(shape)
        if not _shape_intersects(x, y, width, height, header_area):
            continue
        fill_rgb = _shape_fill_rgb(shape)
        text = _shape_text_value(shape)
        if not text and (
            width >= 1.0
            and height >= 0.15
            and (fill_rgb in {"2F1B3A", "2F1838", "321B3D", "3A1D43", "801C80", "811C81", "9D229D"} or y <= 0.25)
        ) or (
            "重庆电子科技职业大学" in text
            or "CHONGQING POLYTECHNIC" in text.upper()
            or text in {"学术汇报", "ACADEMIC REPORT"}
        ):
            if _remove_shape_from_slide(shape):
                removed += 1
    return removed


def _remove_theme_preset_header_text(slide: Any) -> int:
    removed = 0
    header_text_area = (0.0, 0.0, 9.85, 0.85)
    for shape in list(slide.shapes):
        text = _shape_text_value(shape)
        if not text:
            continue
        x, y, width, height = _shape_bounds(shape)
        if _shape_intersects(x, y, width, height, header_text_area):
            if _remove_shape_from_slide(shape):
                removed += 1
    return removed


def _apply_theme_preset_gradient_header(slide: Any, metadata: dict[str, Any]) -> None:
    gradient_from = _theme_color(metadata, "gradient_from", "811C81")
    gradient_to = _theme_color(metadata, "gradient_to", "9D229D")
    _theme_preset_add_gradient_bar(slide, height=0.62, start_color=gradient_from, end_color=gradient_to, slide_width=13.333)


def _apply_theme_preset_content_header(slide: Any, metadata: dict[str, Any], title: str) -> None:
    from pptx.enum.text import PP_ALIGN

    _apply_theme_preset_gradient_header(slide, metadata)
    _add_builtin_template_textbox(
        slide,
        0.58,
        0.15,
        8.8,
        0.34,
        title[:54],
        font_size=16,
        font_family="SimHei",
        color="FFFFFF",
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def _apply_theme_preset_final_slide(slide: Any, metadata: dict[str, Any]) -> None:
    from pptx.enum.text import PP_ALIGN

    for shape in list(slide.shapes):
        _remove_shape_from_slide(shape)
    _theme_preset_add_rect(slide, 0, 0, 13.333, 7.5, "FFFFFF")
    _add_builtin_template_textbox(
        slide,
        3.1,
        2.92,
        7.13,
        0.72,
        "谢谢聆听",
        font_size=34,
        font_family="SimHei",
        color=_theme_color(metadata, "primary", "801C80"),
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_builtin_template_textbox(
        slide,
        4.4,
        3.86,
        4.53,
        0.3,
        "欢迎交流",
        font_size=13,
        font_family="SimHei",
        color="6B556F",
        alignment=PP_ALIGN.CENTER,
    )


def _fix_theme_preset_metric_label_overlap(slide: Any) -> int:
    import re

    from pptx.util import Inches

    changed = 0
    metrics: list[tuple[Any, float, float, float, float]] = []
    labels: list[tuple[Any, float, float, float, float]] = []
    metric_pattern = re.compile(r"^\s*\d+\s*[→\-–—]\s*\d+\s*$")
    for shape in slide.shapes:
        text = _shape_text_value(shape)
        if not text:
            continue
        x, y, width, height = _shape_bounds(shape)
        if metric_pattern.match(text):
            metrics.append((shape, x, y, width, height))
        elif "专业" in text and len(text) <= 12:
            labels.append((shape, x, y, width, height))

    for metric, metric_x, metric_y, metric_width, metric_height in metrics:
        metric_text = _shape_text_value(metric)
        target_metric_width = 1.18 if len(metric_text.replace(" ", "")) <= 4 else 1.42
        metric_right = metric_x + min(metric_width, target_metric_width)
        for label, label_x, label_y, label_width, label_height in labels:
            same_row = abs((label_y + label_height / 2) - (metric_y + metric_height / 2)) <= 0.32
            if not same_row:
                continue
            too_close = label_x < metric_right + 0.16
            overlaps_metric_box = label_x < metric_x + metric_width and label_y < metric_y + metric_height and metric_y < label_y + label_height
            if not (too_close or overlaps_metric_box):
                continue
            try:
                metric.width = Inches(target_metric_width)
                label.left = Inches(metric_x + target_metric_width + 0.18)
                label.width = Inches(max(1.15, 12.78 - (metric_x + target_metric_width + 0.18)))
                changed += 1
            except Exception:
                continue
            break
    return changed


def _write_theme_preset_logo_asset(target_path: Path) -> None:
    try:
        from PIL import Image

        with zipfile.ZipFile(BUILTIN_TEMPLATE_PPTX, "r") as archive:
            with archive.open(BUILTIN_TEMPLATE_LOGO_MEDIA) as image_file:
                image = Image.open(image_file).convert("RGBA")
                pixels = image.load()
                for y in range(image.height):
                    for x in range(image.width):
                        red, green, blue, alpha = pixels[x, y]
                        if alpha > 0 and max(red, green, blue) >= 160:
                            pixels[x, y] = (255, 255, 255, 255)
                        else:
                            pixels[x, y] = (255, 255, 255, 0)
                image.save(target_path, format="PNG")
    except Exception:
        with zipfile.ZipFile(BUILTIN_TEMPLATE_PPTX, "r") as archive:
            target_path.write_bytes(archive.read(BUILTIN_TEMPLATE_LOGO_MEDIA))


def _apply_theme_preset_visual_polish(pptx_path: Path, settings: dict[str, Any], role_mapping: list[dict[str, Any]]) -> int:
    if not _is_builtin_template_request(settings):
        return 0
    metadata = _builtin_template_metadata()
    if _builtin_template_mode(metadata) != "theme_preset":
        return 0
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation(str(pptx_path))
    logo_path = pptx_path.with_name(f"{pptx_path.stem}.image9-logo.tmp.png")
    try:
        _write_theme_preset_logo_asset(logo_path)
    except Exception as exc:
        raise RuntimeError(f"Built-in logo asset is missing: {BUILTIN_TEMPLATE_LOGO_MEDIA}") from exc

    changed = 0
    try:
        slide_count = len(presentation.slides)
        roles_by_slide = {
            int(item.get("slideIndex") or 0): str(item.get("role") or "")
            for item in role_mapping
            if isinstance(item, dict)
        }
        for index, slide in enumerate(presentation.slides, start=1):
            _remove_theme_preset_logo_layer(slide)
            if index == 1:
                _theme_preset_apply_cover(slide, metadata)
                _remove_theme_preset_cover_header(slide)
                changed += 1
                continue
            if index == slide_count or roles_by_slide.get(index) == "ending":
                _apply_theme_preset_final_slide(slide, metadata)
                changed += 1
                continue
            header_title = _theme_preset_title_from_slide(slide, f"第{index}页")
            _remove_theme_preset_top_bar(slide)
            _remove_theme_preset_header_text(slide)
            _apply_theme_preset_content_header(slide, metadata, header_title)
            picture = slide.shapes.add_picture(str(logo_path), Inches(10.35), Inches(0.015), width=Inches(2.45))
            picture.name = "cqet_image9_logo"
            changed += _fix_theme_preset_metric_label_overlap(slide)
            changed += 1
        if changed:
            temp_path = pptx_path.with_name(f"{pptx_path.stem}.theme-brand.tmp{pptx_path.suffix}")
            presentation.save(str(temp_path))
            temp_path.replace(pptx_path)
        return changed
    finally:
        try:
            logo_path.unlink(missing_ok=True)
        except Exception:
            pass


def _slot_alignment(value: str | None) -> Any | None:
    try:
        from pptx.enum.text import PP_ALIGN

        mapping = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        return mapping.get(str(value or "").lower())
    except Exception:
        return None


def _slot_color(slot: dict[str, Any], fallback: str = "3F3F3F") -> str:
    value = str(slot.get("color") or fallback).strip()
    return value[1:] if value.startswith("#") else value


def _slot_font_family(slot: dict[str, Any], fallback: str = "SimHei") -> str:
    value = str(slot.get("fontFamily") or fallback).strip()
    return value or fallback


def _truncate_lines(lines: list[str], max_lines: int | None) -> list[str]:
    if not max_lines or max_lines <= 0:
        return lines
    return lines[:max_lines]


def _fit_text_to_slot(text: str, slot: dict[str, Any]) -> str:
    max_lines = int(slot.get("maxLines") or 0)
    overflow = str(slot.get("overflowPolicy") or "truncate")
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    if not lines and text:
        lines = [_normalized_text_line(text)]
    if overflow in {"truncateList", "split"}:
        lines = _truncate_lines(lines, max_lines)
    elif max_lines == 1 and lines:
        lines = [_normalized_text_line(" ".join(lines))]
    elif max_lines:
        lines = _truncate_lines(lines, max_lines)
    max_chars = int(max(20, float(slot.get("w") or 3.0) * 11))
    fitted: list[str] = []
    for line in lines:
        if len(line) > max_chars:
            fitted.append(line[: max_chars - 1].rstrip() + "…")
        else:
            fitted.append(line)
    return "\n".join(fitted)


def _fill_builtin_template_slot(slide: Any, slot: dict[str, Any], text: str) -> None:
    if not isinstance(slot, dict):
        return
    if not _normalized_text_line(text) and slot.get("overflowPolicy") == "omitIfEmpty":
        return
    from pptx.util import Inches

    x = float(slot.get("x") or 0)
    y = float(slot.get("y") or 0)
    width = float(slot.get("w") or 1)
    height = float(slot.get("h") or 0.4)
    shape = None
    source_shape_index = int(slot.get("sourceShapeIndex") or 0)
    if source_shape_index > 0 and not slot.get("inferred"):
        try:
            candidate = slide.shapes[source_shape_index - 1]
            if hasattr(candidate, "text_frame"):
                shape = candidate
        except Exception:
            shape = None
    if shape is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _set_pptx_shape_text(
        shape,
        _fit_text_to_slot(text, slot),
        font_size=int(slot.get("fontSize") or 14),
        font_family=_slot_font_family(slot),
        color=_slot_color(slot),
        bold=bool(slot.get("bold")),
        alignment=_slot_alignment(slot.get("align")),
    )


def _delete_presentation_slides_except(presentation: Any, keep_slide_numbers: list[int]) -> None:
    slide_id_list = presentation.slides._sldIdLst
    original_slide_ids = list(slide_id_list)
    keep_indexes = [slide_number - 1 for slide_number in keep_slide_numbers]
    keep_index_set = set(keep_indexes)
    for index in reversed(range(len(presentation.slides))):
        if index in keep_index_set:
            continue
        rel_id = slide_id_list[index].rId
        presentation.part.drop_rel(rel_id)
        del slide_id_list[index]
    ordered_slide_ids = [original_slide_ids[index] for index in keep_indexes if 0 <= index < len(original_slide_ids)]
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_id in ordered_slide_ids:
        slide_id_list.append(slide_id)


def _extract_pptx_slide_texts(pptx_path: Path) -> list[list[str]]:
    from pptx import Presentation

    presentation = Presentation(str(pptx_path))
    slides: list[list[str]] = []
    for slide in presentation.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = re.sub(r"\s+", " ", str(shape.text or "")).strip()
            if text:
                texts.append(text)
        slides.append(texts)
    return slides


def _normalized_text_line(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "").replace("\x0b", " ")).strip()


def _clean_generated_text_items(texts: list[str]) -> list[str]:
    skip_exact = {
        "",
        "NexusAI",
        "Academic Presentation",
        "学术报告",
        "学术汇报",
        "汇报封面",
        "报告封面",
        "目录",
        "提纲",
        "Agenda",
        "Contents",
        "Summary",
    }
    result: list[str] = []
    for item in texts:
        text = _normalized_text_line(item)
        if not text or text in skip_exact:
            continue
        if "{{" in text or "}}" in text:
            continue
        if any(prompt in text for prompt in _builtin_powerpoint_default_prompts()):
            continue
        if re.fullmatch(r"\d{1,2}(/\d{1,2})?", text):
            continue
        if text.startswith("数据来源") or text.lower().startswith("source:"):
            continue
        if text not in result:
            result.append(text)
    return result


def _builtin_template_slide_count(pptx_path: Path) -> int:
    from pptx import Presentation

    return len(Presentation(str(pptx_path)).slides)


def _builtin_role_for_generated_slide(texts: list[str], index: int, total: int, target_count: int) -> str:
    combined = " ".join(_clean_generated_text_items(texts))
    lower = combined.lower()
    if index == 1:
        return "cover"
    if target_count >= 4 and index == 2:
        return "toc"
    if index == total:
        return "content_auto" if target_count <= 6 else "summary"
    if any(keyword in lower or keyword in combined for keyword in ("agenda", "contents", "目录", "提纲")):
        return "toc" if index == 2 else "content_fixed"
    has_section_signal = any(
        keyword in lower or keyword in combined
        for keyword in (
            "section divider",
            "chapter divider",
            "chapter transition",
            "part ",
            "chapter ",
            "章节",
            "分章",
            "过渡页",
            "第",
        )
    )
    cleaned = _clean_generated_text_items(texts)
    if has_section_signal and len(cleaned) <= 4 and index not in {1, 2, total}:
        return "section"
    if any(keyword in lower or keyword in combined for keyword in ("chart", "matrix", "timeline", "diagram", "图表", "矩阵", "时间线", "流程")):
        return "chart_auto"
    if any(keyword in lower or keyword in combined for keyword in ("image", "figure", "screenshot", "图片", "图像", "截图", "案例")):
        return "image_text_auto"
    if any(keyword in lower or keyword in combined for keyword in ("compare", "comparison", "对比", "比较")):
        return "content_auto"
    return "content_fixed"


def _builtin_template_hash_seed(settings: dict[str, Any], generated_texts: list[list[str]]) -> int:
    raw = "|".join(
        [
            str(settings.get("taskId") or ""),
            str(settings.get("targetSlides") or ""),
            " ".join(_clean_generated_text_items(generated_texts[0] if generated_texts else [])[:3]),
        ]
    )
    value = 0
    for char in raw:
        value = (value * 131 + ord(char)) % 2_147_483_647
    return value


def _builtin_template_variant_groups(metadata: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    variants = _builtin_template_page_variants(metadata)
    if isinstance(variants, dict):
        return {str(role): [entry for entry in entries if isinstance(entry, dict)] for role, entries in variants.items() if isinstance(entries, list)}
    groups: dict[str, list[dict[str, Any]]] = {}
    for entry in variants:
        if not isinstance(entry, dict):
            continue
        role = str(entry.get("role") or "")
        if role:
            groups.setdefault(role, []).append(entry)
    for entries in groups.values():
        entries.sort(key=lambda item: int(item.get("templateSlideIndex") or item.get("sourceSlideIndex") or 0))
    return groups


def _builtin_template_select_variants(
    metadata: dict[str, Any],
    generated_texts: list[list[str]],
    settings: dict[str, Any],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    target_count = max(1, int(settings.get("targetSlides") or len(generated_texts) or 5))
    total = max(target_count, len(generated_texts), 1)
    groups = _builtin_template_variant_groups(metadata)
    seed = _builtin_template_hash_seed(settings, generated_texts)
    selected: list[dict[str, Any]] = []
    role_mapping: list[dict[str, Any]] = []
    role_use_count: dict[str, int] = {}
    body_role_cycle = ["content_fixed", "content_auto", "chart_auto", "image_text_auto"]
    used_template_slide_indexes: set[int] = set()
    body_slide_count = 0

    def template_slide_index(variant: dict[str, Any] | None) -> int:
        if not isinstance(variant, dict):
            return 0
        try:
            return int(variant.get("templateSlideIndex") or 0)
        except Exception:
            return 0

    def choose_variant(candidate_role: str, *, prefer_unused: bool = True) -> dict[str, Any] | None:
        candidates = groups.get(candidate_role) or []
        if not candidates:
            return None
        role_offset = seed + role_use_count.get(candidate_role, 0)
        ordered = candidates[role_offset % len(candidates) :] + candidates[: role_offset % len(candidates)]
        if prefer_unused:
            for candidate in ordered:
                slide_number = template_slide_index(candidate)
                if slide_number and slide_number not in used_template_slide_indexes:
                    return candidate
        return ordered[0]

    for slide_index in range(1, target_count + 1):
        source_texts = generated_texts[min(slide_index - 1, len(generated_texts) - 1)] if generated_texts else []
        role = _builtin_role_for_generated_slide(source_texts, slide_index, total, target_count)
        reason = "model-intent"
        if slide_index == 1:
            role = "cover"
            reason = "first-slide"
        elif target_count >= 4 and slide_index == 2:
            role = "toc"
            reason = "second-slide-agenda"
        elif slide_index == target_count:
            role = "content_auto" if target_count <= 6 else "summary"
            reason = "last-slide-no-ending"
        elif slide_index == 3 and target_count >= 5:
            role = "section"
            reason = "first-chapter-transition"
        elif role in {"cover", "toc", "section", "cover_alt"}:
            role = "content_fixed"
            reason = "unsafe-role-degraded-to-content"
        elif role in {"content_auto", "image_text_auto", "chart_auto"}:
            body_slide_count += 1
        elif role == "content_fixed" and target_count >= 7:
            role = body_role_cycle[body_slide_count % len(body_role_cycle)]
            body_slide_count += 1
            reason = "body-role-cycle"

        if not groups.get(role) and groups.get("content_fixed"):
            role = "content_fixed"
            reason = f"{reason}-missing-role-fallback"
        variant = choose_variant(role)
        if not variant:
            raise RuntimeError(f"Built-in academic PPT template has no usable variant for role {role}.")
        if role in set(body_role_cycle) and template_slide_index(variant) in used_template_slide_indexes:
            for fallback_role in [role, *[item for item in body_role_cycle if item != role]]:
                fallback_variant = choose_variant(fallback_role)
                if fallback_variant and template_slide_index(fallback_variant) not in used_template_slide_indexes:
                    if fallback_role != role:
                        role = fallback_role
                        reason = f"{reason}-unique-slide-fallback"
                    variant = fallback_variant
                    break
        role_use_count[role] = role_use_count.get(role, 0) + 1
        used_template_slide_indexes.add(template_slide_index(variant))
        selected.append(
            {
                "slideIndex": slide_index,
                "role": role,
                "variantId": variant.get("variantId"),
                "sourceFile": variant.get("sourceFile"),
                "sourceSlideIndex": variant.get("sourceSlideIndex"),
                "templateSlideIndex": variant.get("templateSlideIndex"),
            }
        )
        role_mapping.append(
            {
                "slideIndex": slide_index,
                "role": role,
                "variantId": variant.get("variantId"),
                "reason": reason,
            }
        )
    return selected, role_mapping


def _builtin_variant_by_id(metadata: dict[str, Any], variant_id: str) -> dict[str, Any] | None:
    for entries in _builtin_template_variant_groups(metadata).values():
        for entry in entries:
            if entry.get("variantId") == variant_id:
                return entry
    return None


def _theme_preset_select_variants(
    pptx_path: Path,
    settings: dict[str, Any],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    try:
        actual_slide_count = _builtin_template_slide_count(pptx_path)
    except Exception:
        actual_slide_count = max(1, int(settings.get("targetSlides") or 1))
    generated_texts = _extract_pptx_slide_texts(pptx_path)
    if not generated_texts:
        generated_texts = [[] for _ in range(actual_slide_count)]
    slide_count = max(actual_slide_count, len(generated_texts), 1)
    selected: list[dict[str, Any]] = []
    role_mapping: list[dict[str, Any]] = []
    body_role_cycle = ["content_fixed", "content_auto", "chart_auto", "image_text_auto"]
    body_role_index = 0

    for slide_index in range(1, slide_count + 1):
        source_texts = generated_texts[min(slide_index - 1, len(generated_texts) - 1)] if generated_texts else []
        role = _builtin_role_for_generated_slide(source_texts, slide_index, slide_count, slide_count)
        reason = "model-intent"
        if slide_index == 1:
            role = "cover"
            reason = "first-slide"
        elif slide_count >= 4 and slide_index == 2:
            role = "toc"
            reason = "second-slide-agenda"
        elif slide_count >= 5 and slide_index == 3:
            role = "section"
            reason = "first-chapter-transition"
        elif slide_index == slide_count:
            role = "content_auto" if slide_count <= 6 else "summary"
            reason = "last-slide-no-ending"
        elif role in {"cover", "cover_alt", "toc", "section"}:
            role = body_role_cycle[body_role_index % len(body_role_cycle)]
            body_role_index += 1
            reason = "unsafe-role-degraded-to-content"
        elif role == "content_fixed" and slide_count >= 7:
            role = body_role_cycle[body_role_index % len(body_role_cycle)]
            body_role_index += 1
            reason = "body-role-cycle"
        elif role in body_role_cycle:
            body_role_index += 1

        variant_id = f"theme-{role}-{slide_index:02d}"
        selected.append(
            {
                "slideIndex": slide_index,
                "role": role,
                "variantId": variant_id,
                "sourceFile": "theme-preset",
                "sourceSlideIndex": slide_index,
                "templateSlideIndex": slide_index,
            }
        )
        role_mapping.append(
            {
                "slideIndex": slide_index,
                "role": role,
                "variantId": variant_id,
                "reason": reason,
            }
        )

    return selected, role_mapping


def _clear_builtin_template_slide_placeholders(slide: Any) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        shape_name = str(getattr(shape, "name", "") or "")
        text = str(shape.text or "")
        if _clean_builtin_template_shape_text(text):
            _set_pptx_shape_text(shape, "")
            continue
        if text.strip().isdigit() and (
            "灯片编号" in shape_name
            or "幻灯片编号" in shape_name
            or "Slide Number" in shape_name
            or "SlideNumber" in shape_name
        ):
            _set_pptx_shape_text(shape, "")


def _first_meaningful_text(texts: list[str], *skip_terms: str) -> str:
    for text in texts:
        if not text:
            continue
        if any(term and term in text for term in skip_terms):
            continue
        return text
    return ""


def _section_title_from_generated_texts(texts: list[str]) -> str:
    for text in texts:
        candidate = re.sub(r"\s+", " ", str(text or "")).strip()
        if not candidate or candidate == "学术汇报" or candidate.startswith("数据来源"):
            continue
        candidate = re.sub(r"^第\s*\d+\s*部分[:：\s-]*", "", candidate).strip()
        candidate = re.sub(r"^\d{1,2}\s*", "", candidate).strip()
        if candidate:
            return candidate
    return "章节概览"


def _fill_builtin_template_existing_content_boxes(slide: Any, body: list[str]) -> int:
    if not body:
        return 0
    candidates: list[Any] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or getattr(shape, "is_placeholder", False):
            continue
        shape_name = str(getattr(shape, "name", "") or "")
        if not ("文本框" in shape_name or "TextBox" in shape_name):
            continue
        # The school template body pages already contain empty text boxes inside
        # cards/frames; reuse them instead of drawing a second layout on top.
        if float(getattr(shape, "top", 0) or 0) < 914400:
            continue
        candidates.append(shape)
    candidates.sort(key=lambda item: (int(getattr(item, "left", 0) or 0), int(getattr(item, "top", 0) or 0)))

    filled = 0
    for shape, text in zip(candidates, body):
        if not text:
            continue
        is_heading = filled % 3 == 0
        _set_pptx_shape_text(
            shape,
            text[:80] if is_heading else text[:180],
            font_size=14 if is_heading else 11,
            color="801C80" if is_heading else "3F3F3F",
            bold=is_heading,
        )
        filled += 1
    return filled


def _fill_builtin_template_cover(slide: Any, generated_texts: list[str]) -> None:
    from pptx.enum.text import PP_ALIGN

    title_lines = [text for text in generated_texts if text and "学术汇报" not in text]
    title = "\n".join(title_lines[:2]) or _first_meaningful_text(generated_texts)
    subtitle = title_lines[2] if len(title_lines) > 2 else ""
    author = next((text for text in generated_texts if text.startswith("答辩人") or text.startswith("作者")), "")
    date = next((text for text in reversed(generated_texts) if re.search(r"\d{4}[-/年.]\d{1,2}", text)), "")
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = str(shape.text or "")
        if "{{TITLE}}" in text:
            _set_pptx_shape_text(shape, title, font_size=30, color="3F3F3F", bold=True, alignment=PP_ALIGN.CENTER)
        elif "{{SUBTITLE}}" in text:
            _set_pptx_shape_text(shape, subtitle, font_size=16, color="666666", alignment=PP_ALIGN.CENTER)
        elif "{{AUTHOR}}" in text:
            _set_pptx_shape_text(shape, author, font_size=16, color="3F3F3F", alignment=PP_ALIGN.CENTER)
        elif "{{DATE}}" in text:
            _set_pptx_shape_text(shape, date, font_size=16, color="3F3F3F", bold=True, alignment=PP_ALIGN.CENTER)
        elif "{{SECTION_TITLE}}" in text:
            _set_pptx_shape_text(shape, "学术汇报", font_size=16, color="3F3F3F", bold=True, alignment=PP_ALIGN.CENTER)
        elif "{{SLIDE_TITLE}}" in text:
            _set_pptx_shape_text(shape, "专业分析", font_size=11, color="FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)
        elif "{{BODY}}" in text:
            _set_pptx_shape_text(shape, "建设诊断", font_size=11, color="FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)
        elif "{{KEY_POINTS}}" in text:
            _set_pptx_shape_text(shape, "决策参考", font_size=11, color="FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)


def _fill_builtin_template_toc(slide: Any, generated_texts: list[str]) -> None:
    from pptx.enum.text import PP_ALIGN

    _clear_builtin_template_slide_placeholders(slide)
    items = [
        text
        for text in generated_texts
        if text
        and not text.isdigit()
        and "目录" not in text
        and text != "学术汇报"
        and not text.startswith("数据来源")
        and not re.fullmatch(r"\d{1,2}", text)
    ]
    _add_builtin_template_textbox(slide, 1.25, 1.48, 4.7, 0.5, "目录", font_size=26, color="801C80", bold=True)
    for index, item in enumerate(items[:4], start=1):
        y = 2.08 + (index - 1) * 0.65
        _add_builtin_template_textbox(
            slide,
            1.55,
            y,
            0.35,
            0.3,
            f"{index:02d}",
            font_size=12,
            color="FFFFFF",
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        _add_builtin_template_textbox(slide, 2.0, y - 0.02, 5.8, 0.35, item[:48], font_size=16, color="3F3F3F", bold=True)


def _fill_builtin_template_section(slide: Any, generated_texts: list[str], section_number: int) -> None:
    from pptx.enum.text import PP_ALIGN

    _clear_builtin_template_slide_placeholders(slide)
    title = _section_title_from_generated_texts(generated_texts)
    for shape in slide.shapes:
        if getattr(shape, "name", "") == "文本框 6":
            _set_pptx_shape_text(
                shape,
                f"第{section_number}部分\n{title}",
                font_size=26,
                color="FFFFFF",
                bold=True,
                alignment=PP_ALIGN.LEFT,
            )
            return
    _add_builtin_template_textbox(
        slide,
        5.1,
        2.85,
        6.8,
        1.65,
        f"第{section_number}部分\n{title}",
        font_size=26,
        color="FFFFFF",
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def _fill_builtin_template_content(slide: Any, generated_texts: list[str], title_fallback: str) -> None:
    _clear_builtin_template_slide_placeholders(slide)
    title = _first_meaningful_text(generated_texts, "学术汇报", "数据来源") or title_fallback
    title = re.sub(r"^\d{1,2}\s*", "", title).strip() or title
    body = [
        text
        for text in generated_texts
        if text
        and text != title
        and text != "学术汇报"
        and not text.startswith("数据来源")
        and not re.fullmatch(r"\d{1,2}", text)
    ]
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if int(shape.placeholder_format.type) == 1:
                _set_pptx_shape_text(shape, title[:48], font_size=20, color="801C80", bold=True)
        except Exception:
            continue
    card_texts = body[:12]
    if not card_texts:
        return
    _fill_builtin_template_existing_content_boxes(slide, card_texts)


def _fill_builtin_template_summary(slide: Any, generated_texts: list[str]) -> None:
    from pptx.enum.text import PP_ALIGN

    _clear_builtin_template_slide_placeholders(slide)
    points = [
        text
        for text in generated_texts
        if text
        and text != "学术汇报"
        and not text.startswith("数据来源")
        and not re.fullmatch(r"\d{1,2}", text)
    ]
    title = _first_meaningful_text(points) or "总结"
    _add_builtin_template_textbox(slide, 2.0, 1.35, 9.4, 0.55, title[:40], font_size=26, color="801C80", bold=True, alignment=PP_ALIGN.CENTER)
    for index, item in enumerate(points[1:5], start=1):
        y = 2.05 + (index - 1) * 0.55
        _add_builtin_template_textbox(slide, 2.2, y, 0.35, 0.3, str(index), font_size=12, color="FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)
        _add_builtin_template_textbox(slide, 2.7, y - 0.05, 7.8, 0.42, item[:90], font_size=15, color="3F3F3F")
    closing = points[-1] if points else "谢谢"
    _add_builtin_template_textbox(slide, 2.0, 5.2, 9.4, 0.5, closing[:42], font_size=20, color="801C80", bold=True, alignment=PP_ALIGN.CENTER)


def _builtin_template_keep_slides(target_count: int) -> list[int]:
    if target_count <= 1:
        return [1]
    if target_count == 2:
        return [1, 27]
    if target_count == 3:
        return [1, 3, 27]
    if target_count == 4:
        return [1, 3, 5, 27]
    body_candidates = [5, 6, 7, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 21, 22, 24, 25, 26]
    body_count = max(1, min(target_count - 4, len(body_candidates)))
    return [1, 3, 4, *body_candidates[:body_count], 27]


def _recompose_builtin_template_pptx(pptx_path: Path, settings: dict[str, Any]) -> int:
    if not _is_builtin_template_request(settings):
        return 0
    generated_texts = _extract_pptx_slide_texts(pptx_path)
    if not generated_texts:
        return 0
    from pptx import Presentation

    target_count = int(settings.get("targetSlides") or len(generated_texts) or 5)
    target_count = max(1, target_count)
    keep_slides = _builtin_template_keep_slides(target_count)
    presentation = Presentation(str(BUILTIN_TEMPLATE_PPTX))
    _delete_presentation_slides_except(presentation, keep_slides)
    slides = list(presentation.slides)
    if not slides:
        return 0

    _fill_builtin_template_cover(slides[0], generated_texts[0] if generated_texts else [])
    if len(slides) >= 2:
        _fill_builtin_template_toc(slides[1], generated_texts[1] if len(generated_texts) > 1 else [])
    if len(slides) >= 3 and len(keep_slides) >= 3 and keep_slides[2] == 4:
        _fill_builtin_template_section(slides[2], generated_texts[2] if len(generated_texts) > 2 else [], 1)

    body_start = 3 if len(slides) >= 5 and keep_slides[2] == 4 else 2
    body_end = max(body_start, len(slides) - 1)
    source_body_index = 3 if len(generated_texts) > 4 else 2
    for slide_index in range(body_start, body_end):
        source_texts = generated_texts[min(source_body_index, len(generated_texts) - 1)]
        _fill_builtin_template_content(slides[slide_index], source_texts, f"正文内容 {slide_index - body_start + 1}")
        source_body_index += 1

    if len(slides) >= 2:
        _fill_builtin_template_summary(slides[-1], generated_texts[-1])

    for slide in slides:
        _clear_builtin_template_slide_placeholders(slide)

    temp_path = pptx_path.with_name(f"{pptx_path.stem}.builtin-template.tmp{pptx_path.suffix}")
    presentation.save(str(temp_path))
    temp_path.replace(pptx_path)
    return len(slides)


def _builtin_slide_title(texts: list[str], fallback: str) -> str:
    cleaned = _clean_generated_text_items(texts)
    return (cleaned[0] if cleaned else fallback)[:120]


def _builtin_slide_body(texts: list[str], skip: str = "") -> list[str]:
    cleaned = _clean_generated_text_items(texts)
    return [item for item in cleaned if item != skip][:12]


def _fill_builtin_template_variant(
    slide: Any,
    variant: dict[str, Any],
    role: str,
    generated_texts: list[str],
    slide_index: int,
) -> None:
    _clear_builtin_template_slide_placeholders(slide)
    slots = variant.get("slots") or {}
    cleaned = _clean_generated_text_items(generated_texts)
    title = _builtin_slide_title(generated_texts, "学术汇报")
    body_items = _builtin_slide_body(generated_texts, title)

    def first_slot(slot_type: str) -> dict[str, Any]:
        if isinstance(slots, dict):
            value = slots.get(slot_type)
            return value if isinstance(value, dict) else {}
        if isinstance(slots, list):
            for slot in slots:
                if isinstance(slot, dict) and slot.get("type") == slot_type:
                    return slot
        return {}

    def all_slots(slot_type: str) -> list[dict[str, Any]]:
        if isinstance(slots, dict):
            value = slots.get(slot_type)
            return [value] if isinstance(value, dict) else []
        if isinstance(slots, list):
            return [slot for slot in slots if isinstance(slot, dict) and slot.get("type") == slot_type]
        return []

    if role in {"cover", "cover_alt"}:
        _fill_builtin_template_slot(slide, first_slot("title"), title)
        _fill_builtin_template_slot(slide, first_slot("subtitle"), body_items[0] if body_items else "")
        _fill_builtin_template_slot(slide, first_slot("author"), body_items[1] if len(body_items) > 1 else "")
        _fill_builtin_template_slot(slide, first_slot("organization"), body_items[2] if len(body_items) > 2 else "")
        _fill_builtin_template_slot(slide, first_slot("date"), body_items[3] if len(body_items) > 3 else "")
        return

    if role == "toc":
        items = body_items or cleaned[1:] or cleaned or ["研究背景与问题", "方法与证据", "结论与展望"]
        for index, slot in enumerate(all_slots("tocNumber"), start=1):
            _fill_builtin_template_slot(slide, slot, f"{index:02d}")
        for index, slot in enumerate(all_slots("tocItem")):
            _fill_builtin_template_slot(slide, slot, (items[index] if index < len(items) else "")[:52])
        return

    if role == "section":
        section_number = max(1, slide_index - 2)
        _fill_builtin_template_slot(slide, first_slot("sectionNumber"), f"{section_number:02d}")
        _fill_builtin_template_slot(slide, first_slot("sectionEnglishTitle"), f"SECTION {section_number:02d}")
        _fill_builtin_template_slot(slide, first_slot("sectionTitle") or first_slot("title"), title)
        _fill_builtin_template_slot(slide, first_slot("sectionSubtitle"), body_items[0] if body_items else "")
        return

    if role == "summary":
        _fill_builtin_template_slot(slide, first_slot("slideTitle") or first_slot("title"), title)
        summary_lines = body_items or cleaned[1:] or cleaned[:4]
        for index, slot in enumerate(all_slots("body")):
            _fill_builtin_template_slot(slide, slot, summary_lines[index] if index < len(summary_lines) else "")
        return

    if role == "ending":
        closing = "谢谢" if not cleaned else (cleaned[0] if len(cleaned[0]) <= 12 else "谢谢")
        _fill_builtin_template_slot(slide, first_slot("closingText") or first_slot("title"), closing)
        if len(cleaned) > 1:
            _fill_builtin_template_slot(slide, first_slot("footerNote"), cleaned[1])
        return

    _fill_builtin_template_slot(slide, first_slot("slideTitle") or first_slot("title"), title)
    _fill_builtin_template_slot(slide, first_slot("slideNumber"), f"{slide_index:02d}")
    if first_slot("keyPoints") and body_items:
        _fill_builtin_template_slot(slide, first_slot("keyPoints"), "\n".join(body_items[:4]))
        body_items = body_items[4:] or body_items
    body_source = body_items or cleaned[1:] or cleaned
    body_slots = all_slots("body")
    if body_slots:
        for index, slot in enumerate(body_slots):
            if slot.get("inferred"):
                max_body_lines = int(slot.get("maxLines") or 8)
                _fill_builtin_template_slot(slide, slot, "\n".join(body_source[:max_body_lines]))
                break
            _fill_builtin_template_slot(slide, slot, body_source[index] if index < len(body_source) else "")
    elif variant.get("allowAutoLayout"):
        region = variant.get("autoLayoutRegion") or {}
        if region.get("enabled") is not False:
            auto_slot = {
                **region,
                "type": "body",
                "fontSize": 14,
                "fontFamily": "SimHei",
                "color": "#3F3F3F",
                "align": "left",
                "maxLines": 8,
                "overflowPolicy": "fitThenSplit",
                "inferred": True,
            }
            _fill_builtin_template_slot(slide, auto_slot, "\n".join(body_source[:8]))
    if role == "chart_auto" and all_slots("chart") and body_source:
        for index, slot in enumerate(all_slots("chart")):
            _fill_builtin_template_slot(slide, slot, body_source[index] if index < len(body_source) else "")


def _recompose_builtin_template_pptx(
    pptx_path: Path,
    settings: dict[str, Any],
) -> tuple[int, list[dict[str, Any]], list[dict[str, Any]]]:
    if not _is_builtin_template_request(settings):
        return 0, [], []
    metadata = _builtin_template_metadata()
    generated_texts = _extract_pptx_slide_texts(pptx_path)
    if not generated_texts:
        return 0, [], []
    from pptx import Presentation

    target_count = max(1, int(settings.get("targetSlides") or len(generated_texts) or 5))
    selected_variants, role_mapping = _builtin_template_select_variants(metadata, generated_texts, settings)
    keep_slides = [int(item["templateSlideIndex"]) for item in selected_variants]
    presentation = Presentation(str(BUILTIN_TEMPLATE_PPTX))
    _delete_presentation_slides_except(presentation, keep_slides)
    slides = list(presentation.slides)
    if not slides:
        return 0, selected_variants, role_mapping

    for output_index, (slide, selected) in enumerate(zip(slides, selected_variants), start=1):
        variant = _builtin_variant_by_id(metadata, str(selected.get("variantId") or "")) or {}
        source_texts = generated_texts[min(output_index - 1, len(generated_texts) - 1)]
        _fill_builtin_template_variant(slide, variant, str(selected.get("role") or "content_fixed"), source_texts, output_index)
        _clear_builtin_template_slide_placeholders(slide)

    temp_path = pptx_path.with_name(f"{pptx_path.stem}.builtin-template.tmp{pptx_path.suffix}")
    presentation.save(str(temp_path))
    temp_path.replace(pptx_path)
    return min(len(slides), target_count), selected_variants, role_mapping


def _remove_unfilled_pptx_placeholders(pptx_path: Path, settings: dict[str, Any]) -> int:
    if not _is_builtin_template_request(settings):
        return 0
    return _replace_pptx_text_literals(pptx_path, {token: "" for token in _builtin_placeholder_tokens()})


def _sanitize_powerpoint_default_prompts(pptx_path: Path, settings: dict[str, Any]) -> int:
    if not _is_builtin_template_request(settings):
        return 0
    return _replace_pptx_text_literals(pptx_path, {prompt: "" for prompt in _builtin_powerpoint_default_prompts()})


def _strip_pptx_notes(pptx_path: Path, settings: dict[str, Any]) -> int:
    if not _is_builtin_template_request(settings):
        return 0
    note_relationship_types = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
    )
    removed = 0
    temp_path = pptx_path.with_name(f"{pptx_path.stem}.notes-clean.tmp{pptx_path.suffix}")
    with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            name = item.filename
            if (
                name.startswith("ppt/notesSlides/")
                or name.startswith("ppt/notesMasters/")
                or name.startswith("ppt/notesSlides/_rels/")
                or name.startswith("ppt/notesMasters/_rels/")
            ):
                removed += 1
                continue
            payload = source.read(item)
            if name.endswith(".rels") and payload:
                text = payload.decode("utf-8", errors="ignore")
                original = text
                for rel_type in note_relationship_types:
                    text = re.sub(
                        rf'<Relationship\b(?=[^>]*\bType="{re.escape(rel_type)}")[^>]*/>',
                        "",
                        text,
                    )
                if text != original:
                    removed += 1
                    payload = text.encode("utf-8")
            elif name == "[Content_Types].xml" and payload:
                text = payload.decode("utf-8", errors="ignore")
                original = text
                text = re.sub(r'<Override\b(?=[^>]*\bPartName="/ppt/notesSlides/[^"]+")[^>]*/>', "", text)
                text = re.sub(r'<Override\b(?=[^>]*\bPartName="/ppt/notesMasters/[^"]+")[^>]*/>', "", text)
                if text != original:
                    removed += 1
                    payload = text.encode("utf-8")
            target.writestr(item, payload)
    temp_path.replace(pptx_path)
    return removed


def _openxml_relationships_part(xml_part: str) -> str:
    directory = posixpath.dirname(xml_part)
    filename = posixpath.basename(xml_part)
    if directory:
        return f"{directory}/_rels/{filename}.rels"
    return f"_rels/{filename}.rels"


def _openxml_resolve_relationship_target(relationships_part: str, target: str) -> str:
    normalized_target = str(target or "").replace("\\", "/").strip()
    if not normalized_target:
        return ""
    if normalized_target.startswith("/"):
        return normalized_target.lstrip("/")
    rel_directory = posixpath.dirname(relationships_part)
    owner_directory = "" if rel_directory == "_rels" else posixpath.dirname(rel_directory)
    resolved = posixpath.normpath(posixpath.join(owner_directory, normalized_target))
    return resolved.lstrip("./")


def _is_transient_provider_failure_message(message: str) -> bool:
    lower = str(message or "").lower()
    return any(
        token in lower
        for token in (
            "provider returned html error",
            "provider transient status=520",
            "web server is returning an unknown error",
            "stream interrupted",
            "terminated",
            "und_err_socket",
            "und_err_connect_timeout",
            "fetch failed",
            "network_error",
            "network error",
            "timeout",
            "provider_504",
            "provider_error_502",
            "provider_error_503",
            "provider_error_504",
            "provider_error_520",
            "provider_error_524",
            " 502",
            " 503",
            " 504",
            " 520",
            " 524",
        )
    )


def _should_salvage_repair_output(last_stage: str, raw_output_path: str | None, exc: Exception) -> bool:
    if not raw_output_path:
        return False
    output_path = Path(raw_output_path)
    if not output_path.exists() or not output_path.is_file():
        return False
    message = sanitize_message(exc, 400)
    stage_name = str(last_stage or "").lower()
    repair_signal = "repair" in stage_name or "stage=academic_ppt:repair" in message.lower() or "academic_ppt:repair" in message.lower()
    late_stage = stage_name in {"repair", "postprocess", "export"} or repair_signal
    return late_stage and _is_transient_provider_failure_message(message)


def _sanitize_generated_pptx_openxml(pptx_path: Path) -> dict[str, int]:
    package_rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    notes_rel_types = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
    }
    font_rel_type_tokens = ("font", "embeddedfont")
    removable_prefixes = (
        "ppt/notesSlides/",
        "ppt/notesMasters/",
        "ppt/fonts/",
    )
    removable_files = {"ppt/fontTable.xml"}
    stats = {
        "presentation_nodes_removed": 0,
        "presentation_attributes_removed": 0,
        "relationship_entries_removed": 0,
        "content_type_entries_removed": 0,
        "package_parts_removed": 0,
    }
    temp_path = pptx_path.with_name(f"{pptx_path.stem}.openxml-sanitized.tmp{pptx_path.suffix}")

    def is_removed_part(name: str) -> bool:
        return name in removable_files or any(name.startswith(prefix) for prefix in removable_prefixes)

    with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            name = item.filename
            if is_removed_part(name):
                stats["package_parts_removed"] += 1
                continue

            payload = source.read(item)
            if name == "ppt/presentation.xml":
                root = ET.fromstring(payload)
                for attr_name in ("embedTrueTypeFonts", "saveSubsetFonts"):
                    if root.attrib.pop(attr_name, None) is not None:
                        stats["presentation_attributes_removed"] += 1
                for child_name in ("notesMasterIdLst", "embeddedFontLst"):
                    child = root.find(f"./{{{presentation_ns}}}{child_name}")
                    if child is not None:
                        root.remove(child)
                        stats["presentation_nodes_removed"] += 1
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif name.endswith(".rels"):
                root = ET.fromstring(payload)
                for rel in list(root):
                    if rel.tag != f"{{{package_rel_ns}}}Relationship":
                        continue
                    rel_type = str(rel.get("Type") or "")
                    target_name = _openxml_resolve_relationship_target(name, str(rel.get("Target") or ""))
                    rel_type_lower = rel_type.lower()
                    if (
                        rel_type in notes_rel_types
                        or any(token in rel_type_lower for token in font_rel_type_tokens)
                        or is_removed_part(target_name)
                    ):
                        root.remove(rel)
                        stats["relationship_entries_removed"] += 1
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif name == "[Content_Types].xml":
                root = ET.fromstring(payload)
                for child in list(root):
                    tag = child.tag
                    part_name = str(child.get("PartName") or "")
                    content_type = str(child.get("ContentType") or "").lower()
                    extension = str(child.get("Extension") or "").lower()
                    should_remove = False
                    if tag == f"{{{content_types_ns}}}Override":
                        normalized_part = part_name.lstrip("/")
                        should_remove = is_removed_part(normalized_part) or "font" in content_type
                    elif tag == f"{{{content_types_ns}}}Default":
                        should_remove = extension == "odttf" or "font" in content_type
                    if should_remove:
                        root.remove(child)
                        stats["content_type_entries_removed"] += 1
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(item, payload)

    temp_path.replace(pptx_path)
    return stats


def _validate_generated_pptx_openxml(pptx_path: Path) -> list[str]:
    package_rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    office_rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    rel_attr_names = {
        f"{{{office_rel_ns}}}id",
        f"{{{office_rel_ns}}}embed",
        f"{{{office_rel_ns}}}link",
    }
    errors: list[str] = []

    try:
        with zipfile.ZipFile(pptx_path, "r") as archive:
            names = {name for name in archive.namelist() if not name.endswith("/")}
            relationships_index: dict[str, set[str]] = {}
            xml_roots: dict[str, ET.Element] = {}

            for name in sorted(names):
                if not (name.endswith(".xml") or name.endswith(".rels")):
                    continue
                try:
                    root = ET.fromstring(archive.read(name))
                except Exception as exc:
                    errors.append(f"{name}: XML parse failed ({sanitize_message(exc, 120)})")
                    continue

                if name.endswith(".rels"):
                    rel_ids: set[str] = set()
                    for rel in root.findall(f"{{{package_rel_ns}}}Relationship"):
                        rel_id = str(rel.get("Id") or "").strip()
                        if rel_id:
                            rel_ids.add(rel_id)
                        if str(rel.get("TargetMode") or "").lower() == "external":
                            continue
                        resolved_target = _openxml_resolve_relationship_target(name, str(rel.get("Target") or ""))
                        if not resolved_target:
                            errors.append(f"{name}: empty relationship target")
                        elif resolved_target not in names:
                            errors.append(f"{name}: missing relationship target {resolved_target}")
                    relationships_index[name] = rel_ids
                    continue

                xml_roots[name] = root

            for xml_name, root in xml_roots.items():
                rel_part = _openxml_relationships_part(xml_name)
                rel_ids = relationships_index.get(rel_part, set())
                for element in root.iter():
                    for attr_name, attr_value in element.attrib.items():
                        if attr_name in rel_attr_names and str(attr_value or "").strip() and attr_value not in rel_ids:
                            errors.append(f"{xml_name}: unresolved relationship id {attr_value}")
                if xml_name == "ppt/presentation.xml":
                    if root.find(f"./{{{presentation_ns}}}notesMasterIdLst") is not None:
                        errors.append("ppt/presentation.xml: unexpected notesMasterIdLst remained after sanitization")
                    if root.find(f"./{{{presentation_ns}}}embeddedFontLst") is not None:
                        errors.append("ppt/presentation.xml: unexpected embeddedFontLst remained after sanitization")
                    if root.get("embedTrueTypeFonts") == "1":
                        errors.append("ppt/presentation.xml: embedTrueTypeFonts remained enabled")
                    if root.get("saveSubsetFonts") == "1":
                        errors.append("ppt/presentation.xml: saveSubsetFonts remained enabled")
    except zipfile.BadZipFile as exc:
        return [f"{pptx_path.name}: zip container could not be opened ({sanitize_message(exc, 120)})"]

    return errors


def _extract_slide_count(project_dir: str | None) -> int | None:
    if not project_dir:
        return None
    svg_final = Path(project_dir) / "svg_final"
    if svg_final.exists():
        count = len(list(svg_final.glob("*.svg")))
        return count or None
    return None


def _atomic_write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp = path.with_suffix(f"{path.suffix}.tmp")
    temp.write_text(content, encoding="utf-8")
    temp.replace(path)


def _strict_visual_pipeline_enabled() -> bool:
    return get_settings().academic_ppt_strict_visual_pipeline


def _safe_relative_to_task(task_dir: Path, path: Path) -> str | None:
    try:
        return path.resolve().relative_to(task_dir.resolve()).as_posix()
    except ValueError:
        return None


def _project_file_metadata(task_dir: Path, project_dir: str | None) -> dict[str, Any]:
    if not project_dir:
        return {}
    project_path = Path(project_dir)
    metadata: dict[str, Any] = {}
    for key, file_name in (
        ("manuscript", "manuscript.md"),
        ("designSpec", "design_spec.md"),
        ("criticHistory", "critic_history.json"),
    ):
        file_path = project_path / file_name
        if not file_path.exists() or not file_path.is_file():
            metadata[key] = {"exists": False}
            continue
        metadata[key] = {
            "exists": True,
            "relativePath": _safe_relative_to_task(task_dir, file_path),
            "size": file_path.stat().st_size,
        }
    return metadata


def _svg_final_files(project_dir: str | None) -> list[Path]:
    if not project_dir:
        return []
    svg_final = Path(project_dir) / "svg_final"
    if not svg_final.exists():
        return []
    return sorted(svg_final.glob("*.svg"))


def _write_svg_preview_manifest(task_dir: Path, task_id: str, project_dir: str | None, slide_count: int | None) -> dict[str, Any] | None:
    svg_files = _svg_final_files(project_dir)
    if not svg_files:
        return None
    preview_dir = safe_task_child(task_dir, "previews")
    preview_dir.mkdir(parents=True, exist_ok=True)
    slides: list[dict[str, Any]] = []

    try:
        from backend.generator.svg_finalize.render_ready import prepare_svg_file_for_render
    except Exception:
        prepare_svg_file_for_render = None

    for zero_index, svg_file in enumerate(svg_files):
        page_number = zero_index + 1
        target = preview_dir / f"slide-{page_number:03d}.svg"
        source_for_copy = svg_file
        temp_render_path: Path | None = None
        if prepare_svg_file_for_render:
            try:
                temp_render_path = prepare_svg_file_for_render(svg_file)
                source_for_copy = temp_render_path
            except Exception:
                source_for_copy = svg_file
        try:
            _atomic_write_text(target, source_for_copy.read_text(encoding="utf-8", errors="ignore"))
        finally:
            if temp_render_path is not None:
                temp_render_path.unlink(missing_ok=True)
        slides.append(
            {
                "index": zero_index,
                "pageNumber": page_number,
                "url": f"/api/smart-tools/academic-ppt/tasks/{task_id}/preview/{page_number}",
                "imageUrl": f"/api/smart-tools/academic-ppt/tasks/{task_id}/preview/{page_number}",
                "assetPath": f"previews/slide-{page_number:03d}.svg",
                "source": "svg_final",
                "storageProvider": "local",
                "fileSizeBytes": target.stat().st_size,
            }
        )

    generated_at = __import__("datetime").datetime.utcnow().isoformat() + "Z"
    manifest = {
        "available": True,
        "taskId": task_id,
        "status": "ready",
        "type": "svg",
        "previewType": "svg",
        "source": "svg_final",
        "slideCount": len(slides),
        "previewCount": len(slides),
        "slides": slides,
        "pptxUrl": f"/api/smart-tools/academic-ppt/tasks/{task_id}/download",
        "previewManifestUrl": f"/api/smart-tools/academic-ppt/tasks/{task_id}/preview",
        "previewStoragePrefix": f"academic-ppt/tasks/{task_id}/previews",
        "storageProvider": "local",
        "fallbackReason": None,
        "createdAt": generated_at,
        "updatedAt": generated_at,
        "generatedAt": generated_at,
    }
    _atomic_write_text(preview_dir / "manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2))
    write_checkpoint(
        task_dir,
        "preview",
        {
            "type": "svg",
            "previewType": "svg",
            "slideCount": len(slides),
            "source": "paper-ppt-agent svg_final",
            "targetSlideCount": slide_count,
        },
    )
    return manifest


def _write_visual_pipeline_checkpoints(task_dir: Path, project_dir: str | None, slide_count: int | None) -> None:
    project_metadata = _project_file_metadata(task_dir, project_dir)
    svg_files = _svg_final_files(project_dir)
    composition = _assess_visual_composition_quality(project_dir, slide_count)
    write_checkpoint(
        task_dir,
        "visual-pipeline",
        {
            "project": project_metadata,
            "slideCount": slide_count,
            "svgFinalCount": len(svg_files),
            "contentSlides": composition["contentSlides"],
            "sectionDividerSlides": composition["sectionDividerSlides"],
            "outlineOnlySlides": composition["outlineOnlySlides"],
            "lowContrastTextElements": composition["lowContrastTextElements"],
            "compositionPassed": composition["passed"],
            "compositionIssues": composition["issues"],
            "designSpecComposition": composition["designSpec"],
            "strictVisualPipeline": _strict_visual_pipeline_enabled(),
        },
    )
    write_checkpoint(
        task_dir,
        "svg-final-index",
        {
            "slideCount": len(svg_files),
            "slides": [
                {
                    "index": index,
                    "fileName": svg_file.name,
                    "relativePath": _safe_relative_to_task(task_dir, svg_file),
                    "size": svg_file.stat().st_size,
                }
                for index, svg_file in enumerate(svg_files, start=1)
            ],
        },
    )


_BASIC_FALLBACK_MARKERS = (
    "model bridge was unavailable",
    "this basic deck is generated from parsed text",
    "rule fallback",
    "fallback generated from parsed text",
    "generated a basic structured deck from parsed text",
    "rule fallback -",
)


def _pptx_contains_basic_fallback_markers(output_file: Path) -> bool:
    try:
        with zipfile.ZipFile(output_file) as archive:
            text_parts: list[str] = []
            for name in archive.namelist():
                if not name.startswith("ppt/slides/") or not name.endswith(".xml"):
                    continue
                text_parts.append(archive.read(name).decode("utf-8", errors="ignore"))
        combined = " ".join(text_parts).lower()
        if any(marker in combined for marker in _BASIC_FALLBACK_MARKERS):
            return True
        return "basic deck" in combined and "generated from parsed text" in combined and "model bridge" in combined
    except Exception:
        return False


_BODY_CONTENT_KEYWORDS = (
    "method",
    "evidence",
    "finding",
    "result",
    "analysis",
    "comparison",
    "conclusion",
    "framework",
    "process",
    "chart",
    "figure",
    "table",
    "方法",
    "证据",
    "发现",
    "结果",
    "分析",
    "对比",
    "结论",
    "框架",
    "流程",
    "图",
    "表",
    "工艺",
    "设计",
    "实践",
    "方案",
    "分类",
    "问题",
    "目标",
    "创新",
    "贡献",
    "局限",
)
_AGENDA_KEYWORDS = ("agenda", "outline", "contents", "目录", "提纲", "大纲", "汇报提纲")
_ENDING_KEYWORDS = ("thanks", "thank you", "ending", "结束", "谢谢", "答辩结束")
_DARK_TEXT_ON_DARK_BG = {"#000000", "#020617", "#0b1220", "#0f172a", "#111827", "#172554", "#1e293b"}
_DARK_BG_MARKERS = (
    'fill="#0B1220"',
    "fill='#0B1220'",
    'stop-color="#0B1220"',
    "stop-color='#0B1220'",
    'fill="#172554"',
    "fill='#172554'",
    'stop-color="#172554"',
    "stop-color='#172554'",
    "chapterGradient",
    "dark",
)


def _plain_svg_text(svg: str) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for match in re.finditer(r"<(?:text|tspan)\b([^>]*)>([^<]*)", svg, flags=re.IGNORECASE | re.DOTALL):
        attrs = match.group(1) or ""
        text = re.sub(r"\s+", " ", match.group(2).replace("/text>", "").replace("/tspan>", "")).strip()
        if not text:
            continue
        fill_match = re.search(r"\bfill=[\"']([^\"']+)[\"']", attrs, flags=re.IGNORECASE)
        opacity_match = re.search(r"\bfill-opacity=[\"']([^\"']+)[\"']", attrs, flags=re.IGNORECASE)
        size_match = re.search(r"\bfont-size=[\"']?([0-9.]+)", attrs, flags=re.IGNORECASE)
        x_match = re.search(r"\bx=[\"']?([-0-9.]+)", attrs, flags=re.IGNORECASE)
        y_match = re.search(r"\by=[\"']?([-0-9.]+)", attrs, flags=re.IGNORECASE)
        opacity = 1.0
        if opacity_match:
            try:
                opacity = float(opacity_match.group(1))
            except ValueError:
                opacity = 1.0
        font_size = 16.0
        if size_match:
            try:
                font_size = float(size_match.group(1))
            except ValueError:
                font_size = 16.0
        items.append(
            {
                "text": text,
                "fill": (fill_match.group(1).strip().lower() if fill_match else ""),
                "opacity": opacity,
                "fontSize": font_size,
                "x": float(x_match.group(1)) if x_match else None,
                "y": float(y_match.group(1)) if y_match else None,
                "sourceIndex": match.start(),
            }
        )
    return items


def _svg_has_dark_background(svg: str) -> bool:
    return any(marker.lower() in svg.lower() for marker in _DARK_BG_MARKERS)


def _svg_number(value: str | None, default: float = 0.0) -> float:
    if not value:
        return default
    match = re.search(r"[-+]?\d*\.?\d+", value)
    if not match:
        return default
    try:
        return float(match.group(0))
    except ValueError:
        return default


def _svg_attr(attrs: str, name: str) -> str | None:
    match = re.search(rf"\b{name}=[\"']([^\"']+)[\"']", attrs, flags=re.IGNORECASE)
    return match.group(1).strip() if match else None


def _svg_opacity(attrs: str) -> float:
    return max(0.0, min(1.0, _svg_number(_svg_attr(attrs, "fill-opacity") or _svg_attr(attrs, "opacity"), 1.0)))


def _normalize_svg_color(value: str | None) -> str:
    if not value:
        return ""
    value = value.strip().lower()
    short = re.fullmatch(r"#([0-9a-f])([0-9a-f])([0-9a-f])", value)
    if short:
        return "#" + "".join(channel * 2 for channel in short.groups())
    match = re.fullmatch(r"#[0-9a-f]{6}", value)
    return value if match else value


def _svg_color_luminance(fill: str) -> float | None:
    fill = _normalize_svg_color(fill)
    if not re.fullmatch(r"#[0-9a-f]{6}", fill):
        return None
    channels = [int(fill[index : index + 2], 16) / 255 for index in (1, 3, 5)]
    linear = [value / 12.92 if value <= 0.03928 else ((value + 0.055) / 1.055) ** 2.4 for value in channels]
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]


def _is_dark_svg_fill(fill: str, svg: str) -> bool:
    fill = _normalize_svg_color(fill)
    if fill.startswith("url("):
        return _svg_has_dark_background(svg)
    luminance = _svg_color_luminance(fill)
    return luminance is not None and luminance < 0.18


def _svg_path_bounds(d: str) -> tuple[float, float, float, float] | None:
    tokens = re.findall(r"[A-Za-z]|[-+]?\d*\.?\d+(?:e[-+]?\d+)?", d, flags=re.IGNORECASE)
    if not tokens:
        return None
    command_args = {
        "M": 2,
        "L": 2,
        "T": 2,
        "H": 1,
        "V": 1,
        "C": 6,
        "S": 4,
        "Q": 4,
        "A": 7,
    }
    points: list[tuple[float, float]] = []
    command = ""
    x = 0.0
    y = 0.0
    index = 0
    while index < len(tokens):
        token = tokens[index]
        if re.fullmatch(r"[A-Za-z]", token):
            command = token
            index += 1
            continue
        if not command:
            index += 1
            continue

        upper_command = command.upper()
        arg_count = command_args.get(upper_command)
        if arg_count is None:
            index += 1
            continue
        args: list[float] = []
        lookahead = index
        while lookahead < len(tokens) and len(args) < arg_count and not re.fullmatch(r"[A-Za-z]", tokens[lookahead]):
            args.append(float(tokens[lookahead]))
            lookahead += 1
        if len(args) < arg_count:
            index = lookahead + 1 if lookahead == index else lookahead
            continue

        relative = command.islower()
        if upper_command in {"M", "L", "T"}:
            x = x + args[-2] if relative else args[-2]
            y = y + args[-1] if relative else args[-1]
            points.append((x, y))
        elif upper_command == "H":
            x = x + args[0] if relative else args[0]
            points.append((x, y))
        elif upper_command == "V":
            y = y + args[0] if relative else args[0]
            points.append((x, y))
        elif upper_command in {"C", "S", "Q", "A"}:
            x = x + args[-2] if relative else args[-2]
            y = y + args[-1] if relative else args[-1]
            points.append((x, y))
        if upper_command == "M":
            command = "l" if relative else "L"
        index += arg_count
    if not points:
        return None
    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    return min(xs), min(ys), max(xs), max(ys)


def _svg_background_shapes(svg: str) -> list[dict[str, Any]]:
    shapes: list[dict[str, Any]] = []
    for match in re.finditer(r"<(rect|path)\b([^>]*)>", svg, flags=re.IGNORECASE | re.DOTALL):
        tag = match.group(1).lower()
        attrs = match.group(2) or ""
        fill = _normalize_svg_color(_svg_attr(attrs, "fill"))
        if not fill or fill in {"none", "transparent"}:
            continue
        bounds: tuple[float, float, float, float] | None = None
        if tag == "rect":
            x = _svg_number(_svg_attr(attrs, "x"), 0.0)
            y = _svg_number(_svg_attr(attrs, "y"), 0.0)
            width = _svg_number(_svg_attr(attrs, "width"), 0.0)
            height = _svg_number(_svg_attr(attrs, "height"), 0.0)
            if width > 0 and height > 0:
                bounds = (x, y, x + width, y + height)
        elif tag == "path":
            bounds = _svg_path_bounds(_svg_attr(attrs, "d") or "")
        if not bounds:
            continue
        shapes.append({"sourceIndex": match.start(), "bounds": bounds, "fill": fill, "opacity": _svg_opacity(attrs)})
    return shapes


def _svg_text_background_is_dark(svg: str, item: dict[str, Any], shapes: list[dict[str, Any]]) -> bool:
    x = item.get("x")
    y = item.get("y")
    if not isinstance(x, (int, float)) or not isinstance(y, (int, float)):
        return _svg_has_dark_background(svg)
    for shape in reversed(shapes):
        if shape["sourceIndex"] > item["sourceIndex"] or shape["opacity"] < 0.25:
            continue
        left, top, right, bottom = shape["bounds"]
        if left <= x <= right and top <= y <= bottom:
            return _is_dark_svg_fill(shape["fill"], svg)
    return _svg_has_dark_background(svg)


def _low_contrast_text_count(svg: str, text_items: list[dict[str, Any]]) -> int:
    if not _svg_has_dark_background(svg):
        return 0
    shapes = _svg_background_shapes(svg)
    count = 0
    for item in text_items:
        fill = item["fill"]
        if (
            fill in _DARK_TEXT_ON_DARK_BG
            and item["opacity"] >= 0.45
            and item["fontSize"] >= 20
            and _svg_text_background_is_dark(svg, item, shapes)
        ):
            count += 1
    return count


def _classify_svg_slide(svg_file: Path, index: int, total: int) -> dict[str, Any]:
    svg = svg_file.read_text(encoding="utf-8", errors="ignore")
    text_items = _plain_svg_text(svg)
    texts = [item["text"] for item in text_items]
    combined = " ".join(texts)
    lower = combined.lower()
    meaningful_text_count = len([text for text in texts if len(text.strip()) >= 2])
    has_agenda = any(keyword in lower or keyword in combined for keyword in _AGENDA_KEYWORDS)
    has_ending = index == total or any(keyword in lower or keyword in combined for keyword in _ENDING_KEYWORDS)
    has_part = bool(re.search(r"\bpart\s*\d+\b|\bchapter\s*\d+\b|\bsection\s*\d+\b|章节|篇章", combined, flags=re.IGNORECASE))
    has_body_keyword = any(keyword in lower or keyword in combined for keyword in _BODY_CONTENT_KEYWORDS)
    is_cover = index == 1
    is_content = (not is_cover and not has_agenda and not has_ending and not has_part) and (
        meaningful_text_count >= 5 or (meaningful_text_count >= 4 and has_body_keyword)
    )
    is_section = (not is_cover and not has_agenda and not has_ending and not is_content) and (
        has_part or meaningful_text_count <= 4
    )
    is_outline_only = (has_agenda or (not is_cover and not has_ending and not is_content and meaningful_text_count <= 4))
    return {
        "index": index,
        "fileName": svg_file.name,
        "category": (
            "cover"
            if is_cover
            else "agenda"
            if has_agenda
            else "ending"
            if has_ending
            else "content"
            if is_content
            else "section"
            if is_section
            else "outline_only"
        ),
        "textCount": meaningful_text_count,
        "hasBodyKeyword": has_body_keyword,
        "lowContrastTextCount": _low_contrast_text_count(svg, text_items),
    }


def _assess_design_spec_composition(project_dir: str | None, target_slide_count: int | None) -> dict[str, Any]:
    if not project_dir:
        return {"available": False}
    design_spec = Path(project_dir) / "design_spec.md"
    if not design_spec.exists():
        return {"available": False}
    text = design_spec.read_text(encoding="utf-8", errors="ignore")
    slide_blocks = re.split(r"(?=####\s+Slide\s+\d+|###\s+Slide\s+\d+)", text, flags=re.IGNORECASE)
    slide_blocks = [block for block in slide_blocks if re.search(r"\bSlide\s+\d+\b", block, flags=re.IGNORECASE)]
    content_slides = 0
    section_divider_slides = 0
    outline_only_slides = 0
    for block in slide_blocks:
        block_lower = block.lower()
        is_cover = "cover" in block_lower
        is_agenda = any(keyword in block_lower or keyword in block for keyword in _AGENDA_KEYWORDS)
        is_ending = any(keyword in block_lower or keyword in block for keyword in _ENDING_KEYWORDS)
        is_section = "section" in block_lower or "chapter" in block_lower or "章节" in block
        has_content = bool(re.search(r"\*\*Content\*\*|Content:|内容|Visualization|可视化", block, flags=re.IGNORECASE))
        if not is_cover and not is_agenda and not is_ending and has_content and not is_section:
            content_slides += 1
        if not is_cover and not is_agenda and not is_ending and is_section:
            section_divider_slides += 1
        if is_agenda or (is_section and not has_content):
            outline_only_slides += 1
    return {
        "available": True,
        "totalSlides": len(slide_blocks) or target_slide_count,
        "contentSlides": content_slides,
        "sectionDividerSlides": section_divider_slides,
        "outlineOnlySlides": outline_only_slides,
    }


def _assess_visual_composition_quality(project_dir: str | None, target_slide_count: int | None) -> dict[str, Any]:
    svg_files = _svg_final_files(project_dir)
    total = len(svg_files)
    slide_assessments = [
        _classify_svg_slide(svg_file, index, total)
        for index, svg_file in enumerate(svg_files, start=1)
    ]
    content_slides = sum(1 for item in slide_assessments if item["category"] == "content")
    section_divider_slides = sum(1 for item in slide_assessments if item["category"] == "section")
    outline_only_slides = sum(1 for item in slide_assessments if item["category"] in {"agenda", "outline_only"})
    low_contrast_text = sum(int(item["lowContrastTextCount"]) for item in slide_assessments)
    requested = target_slide_count or total
    min_content = max(1, math.ceil(requested * 0.55)) if requested >= 8 else max(1, requested - 3)
    max_sections = max(1, math.floor(requested * 0.30)) if requested >= 8 else 2
    max_outline = max(1, math.floor(requested * 0.35)) if requested >= 8 else 2
    issues: list[str] = []
    if requested >= 8 and content_slides < min_content:
        issues.append(f"Only {content_slides} body content slides were detected; expected at least {min_content}.")
    if requested >= 8 and section_divider_slides > max_sections:
        issues.append(f"{section_divider_slides} section divider slides were detected; expected no more than {max_sections}.")
    if requested >= 8 and outline_only_slides > max_outline:
        issues.append(f"{outline_only_slides} outline-only slides were detected; expected no more than {max_outline}.")
    if low_contrast_text:
        issues.append(f"{low_contrast_text} dark-template text element(s) appear low contrast on dark backgrounds.")
    if target_slide_count and total > target_slide_count + max(1, math.floor(target_slide_count * 0.15)):
        issues.append(f"{total} SVG pages were generated for a {target_slide_count}-page request.")
    design_spec = _assess_design_spec_composition(project_dir, target_slide_count)
    if requested >= 8 and design_spec.get("available"):
        design_content = int(design_spec.get("contentSlides") or 0)
        design_sections = int(design_spec.get("sectionDividerSlides") or 0)
        design_outline = int(design_spec.get("outlineOnlySlides") or 0)
        if design_content < min_content:
            issues.append(
                f"Design spec planned only {design_content} body content slides; expected at least {min_content}."
            )
        if design_sections > max_sections:
            issues.append(
                f"Design spec planned {design_sections} section divider slides; expected no more than {max_sections}."
            )
        if design_outline > max_outline:
            issues.append(
                f"Design spec planned {design_outline} outline-only slides; expected no more than {max_outline}."
            )
    return {
        "contentSlides": content_slides,
        "sectionDividerSlides": section_divider_slides,
        "outlineOnlySlides": outline_only_slides,
        "lowContrastTextElements": low_contrast_text,
        "minContentSlides": min_content,
        "maxSectionDividerSlides": max_sections,
        "maxOutlineOnlySlides": max_outline,
        "slides": slide_assessments,
        "issues": issues,
        "passed": not issues,
        "designSpec": design_spec,
    }


def _assess_builtin_template_role_plan(
    output_file: Path,
    target_slide_count: int | None,
    settings: dict[str, Any] | None,
    selected_variants: list[dict[str, Any]] | None,
    role_mapping: list[dict[str, Any]] | None,
) -> dict[str, Any]:
    if not settings or not _is_builtin_template_request(settings):
        return {"available": False, "passed": False, "issues": ["not a built-in template request"]}
    if not selected_variants or not role_mapping:
        return {"available": False, "passed": False, "issues": ["missing built-in template role plan"]}

    expected_count = max(1, int(target_slide_count or len(selected_variants)))
    issues: list[str] = []
    roles = [str(item.get("role") or "") for item in selected_variants]
    slide_indexes = [int(item.get("slideIndex") or 0) for item in selected_variants]
    role_ids = [str(item.get("variantId") or "") for item in selected_variants]
    allowed_body_roles = {"content_fixed", "content_auto", "image_text_auto", "chart_auto"}

    if len(selected_variants) != expected_count:
        issues.append(f"Built-in template selected {len(selected_variants)} variants; expected {expected_count}.")
    if role_mapping and len(role_mapping) != len(selected_variants):
        issues.append("Built-in template roleMapping length does not match selectedVariants.")
    try:
        pptx_slide_count = _builtin_template_slide_count(output_file)
        if pptx_slide_count != expected_count:
            issues.append(f"Built-in template PPTX has {pptx_slide_count} slides; expected {expected_count}.")
    except Exception as exc:
        issues.append(f"Built-in template PPTX could not be reopened: {sanitize_message(exc, 160)}")

    for expected_index, slide_index in enumerate(slide_indexes, start=1):
        if slide_index != expected_index:
            issues.append("Built-in template slide indexes are not sequential.")
            break
    if any(not variant_id for variant_id in role_ids):
        issues.append("Built-in template selected a variant without variantId.")
    if roles and roles[0] != "cover":
        issues.append("Built-in template first slide is not cover.")
    if expected_count >= 4 and len(roles) >= 2 and roles[1] != "toc":
        issues.append("Built-in template second slide is not toc.")
    if expected_count >= 5 and len(roles) >= 3 and roles[2] != "section":
        issues.append("Built-in template third slide is not section.")
    if expected_count >= 3 and roles and roles[-1] not in {"summary", "ending"}:
        issues.append("Built-in template last slide is not summary/ending.")

    interior_roles = roles[3:-1] if expected_count >= 5 else roles[2:-1]
    unsafe_interior = [role for role in interior_roles if role not in allowed_body_roles]
    if unsafe_interior:
        issues.append(f"Built-in template body slides used unsafe role(s): {', '.join(sorted(set(unsafe_interior)))}.")
    if expected_count >= 6 and not any(role in allowed_body_roles for role in interior_roles):
        issues.append("Built-in template has no body content variants between section and final slide.")
    if any(role in {"cover", "toc"} for role in roles[2:]):
        issues.append("Built-in template reused cover/toc outside reserved positions.")
    if "section" in roles[3:]:
        issues.append("Built-in template reused section as a body slide.")

    return {
        "available": True,
        "passed": not issues,
        "issues": issues,
        "roles": roles,
        "selectedVariantCount": len(selected_variants),
        "expectedSlideCount": expected_count,
    }


def _visual_pipeline_status(
    project_dir: str | None,
    output_file: Path,
    target_slide_count: int | None = None,
    settings: dict[str, Any] | None = None,
    selected_variants: list[dict[str, Any]] | None = None,
    role_mapping: list[dict[str, Any]] | None = None,
) -> tuple[str, str | None, bool]:
    if _pptx_contains_basic_fallback_markers(output_file):
        return "degraded", "Model bridge unavailable; generated basic parsed-text deck.", True
    if not project_dir:
        return "degraded", "Visual pipeline project directory was not reported by paper-ppt-agent.", False

    project_path = Path(project_dir)
    manuscript = project_path / "manuscript.md"
    design_spec = project_path / "design_spec.md"
    svg_final = project_path / "svg_final"
    svg_files = list(svg_final.glob("*.svg")) if svg_final.exists() else []
    if not manuscript.exists() or manuscript.stat().st_size < 200:
        return "degraded", "Research manuscript was not generated.", False
    if not design_spec.exists() or design_spec.stat().st_size < 200:
        return "degraded", "Strategy design spec was not generated.", False
    if not svg_files:
        return "degraded", "SVG executor output was not generated.", False
    if target_slide_count and len(svg_files) < max(1, target_slide_count - 1):
        return "degraded", "SVG executor output is incomplete for the requested slide count.", False
    composition = _assess_visual_composition_quality(project_dir, target_slide_count)
    if not composition.get("passed", True):
        if settings and _builtin_template_uses_layout_blueprint(settings):
            builtin_role_plan = _assess_builtin_template_role_plan(
                output_file,
                target_slide_count,
                settings,
                selected_variants,
                role_mapping,
            )
            if builtin_role_plan.get("passed"):
                return "success", None, False
        return "degraded", "Visual composition quality gate failed: " + "; ".join(composition.get("issues") or []), False
    return "success", None, False


def _read_model_bridge_snapshot(task_dir: Path) -> dict[str, Any]:
    task_json = task_dir / "task.json"
    try:
        data = json.loads(task_json.read_text(encoding="utf-8"))
    except Exception:
        return {}
    result: dict[str, Any] = {}
    for key in (
        "modelBridgeStatus",
        "modelBridgePrimaryModel",
        "modelBridgePrimaryStatus",
        "modelBridgeFallbackModel",
        "modelBridgeFallbackStatus",
        "modelBridgeErrorSummary",
    ):
        value = data.get(key)
        if value is not None:
            result[key] = value
    return result


def _extract_plain_text_from_input(input_path: Path) -> str:
    suffix = input_path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown", ".tex"}:
        return input_path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            presentation = Presentation(str(input_path))
            chunks: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                if texts:
                    chunks.append(f"Slide {index}\n" + "\n".join(texts))
            return "\n\n".join(chunks)
        except Exception:
            return input_path.stem
    if suffix == ".pdf":
        try:
            import fitz

            with fitz.open(str(input_path)) as doc:
                return "\n\n".join(page.get_text("text") for page in doc[:8])
        except Exception:
            return input_path.stem
    return input_path.stem


def _split_text_units(text: str) -> list[str]:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if not cleaned:
        return []
    chunks = [item.strip() for item in re.split(r"(?<=[。！？.!?])\s+", cleaned) if item.strip()]
    if len(chunks) >= 4:
        return chunks
    return [item.strip() for item in textwrap.wrap(cleaned, width=90) if item.strip()]


def _rule_fallback_slides(input_path: Path, settings: dict[str, Any], text: str) -> list[tuple[str, list[str]]]:
    target = max(4, min(int(settings.get("targetSlides") or 5), 8))
    units = _split_text_units(text)
    title = input_path.stem.replace("_", " ").strip() or "Academic Presentation"
    agenda = ["Research background", "Core problem", "Key evidence", "Summary and next steps"]
    summary = [
        "Model bridge was unavailable.",
        "This basic deck is generated from parsed text.",
        "Please retry for model-enhanced design when the model service is restored.",
    ]

    slides: list[tuple[str, list[str]]] = [
        (title, [units[0][:120] if units else "Generated from parsed source text."]),
        ("Agenda", agenda),
    ]
    content_titles = agenda[: max(target - 3, 1)]
    cursor = 1
    for heading in content_titles:
        bullets = [unit[:140] for unit in units[cursor : cursor + 4]]
        cursor += 4
        if not bullets:
            bullets = ["Content is summarized from the uploaded source."]
        slides.append((heading, bullets))
    slides.append(("Summary", summary))
    return slides[:target]


def _write_rule_fallback_pptx(task_dir: Path, input_path: Path, settings: dict[str, Any], reason: str) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    output_path = safe_task_child(task_dir, "outputs", "academic-ppt-result.pptx")
    paper_workspace = safe_task_child(task_dir, "checkpoints", "paper-workspaces", "rule-fallback")
    svg_output = safe_task_child(paper_workspace, "svg_output")
    svg_final = safe_task_child(paper_workspace, "svg_final")
    exports = safe_task_child(paper_workspace, "exports")
    for directory in (svg_output, svg_final, exports):
        directory.mkdir(parents=True, exist_ok=True)

    text = _extract_plain_text_from_input(input_path)
    slides = _rule_fallback_slides(input_path, settings, text)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for index, (title, bullets) in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 250, 252)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title[:80]
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

        accent = slide.shapes.add_shape(1, Inches(0.7), Inches(1.45), Inches(1.1), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
        accent.line.fill.background()

        body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(11.4), Inches(4.7))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()
        for bullet_index, bullet in enumerate(bullets[:5]):
            paragraph = body_frame.paragraphs[0] if bullet_index == 0 else body_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 65, 85)
            paragraph.space_after = Pt(10)

        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.25))
        footer = footer_box.text_frame.paragraphs[0]
        footer.text = f"Rule fallback - {index}/{len(slides)}"
        footer.alignment = PP_ALIGN.RIGHT
        footer.font.size = Pt(9)
        footer.font.color.rgb = RGBColor(100, 116, 139)

        svg_stub = f"<svg xmlns='http://www.w3.org/2000/svg' width='1280' height='720'><text x='64' y='96'>{title}</text></svg>"
        (svg_output / f"slide_{index:02d}.svg").write_text(svg_stub, encoding="utf-8")
        (svg_final / f"slide_{index:02d}.svg").write_text(svg_stub, encoding="utf-8")

    presentation.save(output_path)
    export_path = exports / "presentation_rule_fallback.pptx"
    shutil.copy2(output_path, export_path)
    design_spec = paper_workspace / "design_spec.md"
    design_spec.write_text(
        "# Rule fallback design spec\n\nModel bridge unavailable; generated a basic structured deck from parsed text.\n",
        encoding="utf-8",
    )
    output_size = output_path.stat().st_size
    write_checkpoint(
        task_dir,
        "pptx-exported",
        {
            "outputFileName": output_path.name,
            "outputFileSize": output_size,
            "slideCount": len(slides),
            "generator": "paper-ppt-agent",
            "generationMode": "paper-ppt-agent-rule-fallback",
            "visualPipelineStatus": "degraded",
            "fallbackReason": reason,
        },
    )
    append_log(task_dir, "warn", "Model enhancement is unavailable; generated a basic PPT from parsed text.")
    return {
        "outputFileName": output_path.name,
        "outputFileSize": output_size,
        "slideCount": len(slides),
        "generationMode": "paper-ppt-agent-rule-fallback",
        "visualPipelineStatus": "degraded",
        "fallbackReason": reason,
        "modelBridgeStatus": "failed",
        "modelBridgePrimaryModel": "subrouter:gpt-5.4",
        "modelBridgePrimaryStatus": "failed",
        "modelBridgeFallbackModel": "moonshot:kimi-k2.5",
        "modelBridgeFallbackStatus": "failed",
        "modelBridgeErrorSummary": reason,
    }

def _is_model_bridge_failure(exc: Exception) -> bool:
    message = sanitize_message(exc).lower()
    return (
        "model bridge" in message
        or "model provider" in message
        or "model unavailable" in message
        or "fallback model unavailable" in message
        or "primary model unavailable" in message
        or "gpt-5.4 visual" in message
        or "strict visual" in message
    )


async def run_paper_ppt_agent_task(
    *,
    task_id: str,
    task_dir: Path,
    input_file_path: Path,
    settings: dict[str, Any],
    request_options: dict[str, Any],
    model_bridge_url: str | None,
    resume: bool,
    resume_from_step: str | None,
    on_progress,
    cancel_event: asyncio.Event,
) -> dict[str, Any]:
    _patch_paper_ppt_agent_runtime(task_dir, model_bridge_url)
    from backend.orchestrator.pipeline import GenerationRequest, run_pipeline

    normalized_input, source_type = _normalize_input_for_paper_agent(input_file_path, task_dir)
    research_context = await _prepare_research_context(
        task_dir=task_dir,
        input_path=input_file_path,
        settings=settings,
        request_options=request_options,
    )
    append_log(task_dir, "info", "Prepared generation input and started the academic PPT engine.")
    write_checkpoint(
        task_dir,
        "generation-state",
        {
            "taskId": task_id,
            "status": "running",
            "sourceType": source_type,
            "resume": resume,
            "resumeFromStep": resume_from_step,
            "deepResearchEnabled": _request_flag(settings, request_options, "deepResearchEnabled", "enableDeepResearch"),
            "externalResearchEnabled": _request_flag(settings, request_options, "externalResearchEnabled", "enableExternalResearch"),
            "webSearchEnabled": bool(request_options.get("webSearchEnabled") or settings.get("webSearchEnabled")),
            "searchProvider": "nexus-searxng",
            "searchStatus": research_context["searchStatus"],
        },
    )

    request = GenerationRequest(
        file_path=normalized_input,
        source_type=source_type,
        provider="nexus",
        model="nexus-primary",
        api_key=task_id,
        base_url=model_bridge_url,
        canvas_format=_canvas_format(settings),
        style=_style(settings),
        num_pages=int(settings.get("targetSlides") or 12),
        instruction="\n\n".join(part for part in [_instruction(settings), research_context["instruction"]] if part),
        language="en" if settings.get("outputLanguage") == "en" else "zh",
        detail_level=_detail_level(settings),
        style_overrides=_visual_quality_style_overrides(settings),
        timeout_seconds=int(get_settings().task_max_runtime_seconds),
        enable_visual_critic=bool(settings.get("enableVisualQa")),
        enable_icon=bool(settings.get("enableIconDecoration")),
        enable_icon_rag=False,
        template_id=_template_id(settings),
    )
    setattr(request, "job_id", task_id)

    output_path: str | None = None
    project_dir: str | None = None
    last_progress = 0
    last_stage = "starting"
    repair_warning_reason: str | None = None

    try:
        async for event in run_pipeline(request):
            if cancel_event.is_set():
                raise asyncio.CancelledError()
            progress = int(max(0, min(100, float(getattr(event, "progress", 0) or 0) * 100)))
            if progress < last_progress:
                progress = last_progress
            last_progress = progress

            stage = str(getattr(event, "stage", "generation"))
            last_stage = stage
            status = str(getattr(event, "status", "progress"))
            message = str(getattr(event, "message", "") or stage)
            data = getattr(event, "data", None) or {}
            if isinstance(data, dict):
                if data.get("project_dir"):
                    project_dir = str(data["project_dir"])
                if data.get("output_path"):
                    output_path = str(data["output_path"])
                if data.get("svg") and data.get("page"):
                    write_checkpoint(
                        task_dir,
                        f"slide-svg-{int(data['page']):03d}",
                        {
                            "page": int(data["page"]),
                            "stage": stage,
                            "status": status,
                            "hasSvg": True,
                        },
                    )

            append_log(task_dir, "error" if status == "error" else "info", _product_log_message(message))
            if project_dir and stage in {"research", "strategy", "generation", "postprocess", "export"}:
                _write_visual_pipeline_checkpoints(task_dir, project_dir, int(settings.get("targetSlides") or 0) or None)
            write_checkpoint(
                task_dir,
                "generation-state",
                {
                    "taskId": task_id,
                    "status": status,
                    "stage": stage,
                    "progress": progress,
                    "message": message,
                    "projectDirKnown": bool(project_dir),
                    "outputKnown": bool(output_path),
                },
            )
            await on_progress(stage, progress, message)
    except Exception as exc:
        if not _is_model_bridge_failure(exc):
            raise
        if _should_salvage_repair_output(last_stage, output_path, exc):
            repair_warning_reason = (
                "completed_with_repair_warning: transient repair-stage provider failure; "
                f"preserved the generated PPTX and will validate it before delivery. {sanitize_message(exc, 180)}"
            )
            append_log(task_dir, "warn", "Repair-stage model bridge failure detected after PPTX generation.")
            append_log(task_dir, "warn", repair_warning_reason)
            await on_progress("export", max(last_progress, 92), "Repair was interrupted; validating generated PPTX.")
        elif _strict_visual_pipeline_enabled():
            strict_reason = (
                "GPT-5.4 visual generation stage is temporarily unavailable; "
                "strict visual mode did not generate a low-quality fallback deck. "
                "Please retry or resume the task later."
            )
            append_log(task_dir, "warn", "Model bridge final failed.")
            append_log(task_dir, "warn", strict_reason)
            append_log(task_dir, "error", "Strict visual pipeline is enabled; not generating rule fallback PPTX.")
            write_checkpoint(
                task_dir,
                "generation-state",
                {
                    "taskId": task_id,
                    "status": "failed",
                    "resumable": True,
                    "strictVisualPipeline": True,
                    "fallbackPrevented": True,
                    "fallbackReason": strict_reason,
                    "sourceReason": sanitize_message(exc, 180),
                },
            )
            raise RuntimeError(strict_reason) from exc
        else:
            reason = f"Model bridge unavailable; generated from parsed text with rule fallback. {sanitize_message(exc, 180)}"
            append_log(task_dir, "warn", "Model bridge final failed.")
            append_log(task_dir, "warn", reason)
            await on_progress("export", 85, "Generating basic PPTX from parsed text.")
            fallback_result = _write_rule_fallback_pptx(task_dir, input_file_path, settings, reason)
            fallback_result.update(
                {
                    "searchStatus": research_context["searchStatus"],
                    "researchStatus": research_context["researchStatus"],
                    "researchSourcesCount": research_context["researchSourcesCount"],
                    "researchFallbackReason": research_context["researchFallbackReason"],
                }
            )
            return fallback_result

    slide_count = _extract_slide_count(project_dir)
    target_slide_count = int(settings.get("targetSlides") or request_options.get("slideCount") or 0) or None
    builtin_text_cleaned = _sanitize_builtin_template_svg_text(project_dir, settings)
    output_file, output_size = _safe_output_file(task_dir, output_path)
    selected_variants: list[dict[str, Any]] = []
    role_mapping: list[dict[str, Any]] = []
    if _is_builtin_template_request(settings):
        pptx_recomposed_slides = 0
        if _builtin_template_uses_layout_blueprint(settings):
            pptx_recomposed_slides, selected_variants, role_mapping = _recompose_builtin_template_pptx(output_file, settings)
        else:
            selected_variants, role_mapping = _theme_preset_select_variants(output_file, settings)
        if pptx_recomposed_slides:
            slide_count = pptx_recomposed_slides
        theme_preset_branded_slides = _apply_theme_preset_visual_polish(output_file, settings, role_mapping)
        pptx_text_cleaned = _replace_pptx_text_literals(output_file, _builtin_generated_text_replacements(settings))
        pptx_placeholder_cleaned = _remove_unfilled_pptx_placeholders(output_file, settings)
        pptx_prompt_cleaned = _sanitize_powerpoint_default_prompts(output_file, settings)
        pptx_notes_cleaned = _strip_pptx_notes(output_file, settings)
        openxml_sanitized = _sanitize_generated_pptx_openxml(output_file)
        openxml_errors = _validate_generated_pptx_openxml(output_file)
        if openxml_errors:
            raise RuntimeError("OpenXML validation failed after sanitization: " + "; ".join(openxml_errors[:8]))
        if (
            builtin_text_cleaned
            or pptx_recomposed_slides
            or theme_preset_branded_slides
            or pptx_text_cleaned
            or pptx_placeholder_cleaned
            or pptx_prompt_cleaned
            or pptx_notes_cleaned
            or any(openxml_sanitized.values())
        ):
            append_log(
                task_dir,
                "info",
                (
                    f"Built-in template cleanup applied: SVG text cleanup {builtin_text_cleaned}, "
                    f"PPTX template recomposed slides {pptx_recomposed_slides}, "
                    f"PPTX theme preset branded slides {theme_preset_branded_slides}, "
                    f"PPTX text cleanup {pptx_text_cleaned}, "
                    f"PPTX placeholder cleanup {pptx_placeholder_cleaned}, "
                    f"PPTX PowerPoint prompt cleanup {pptx_prompt_cleaned}, "
                    f"PPTX notes cleanup {pptx_notes_cleaned}, "
                    f"OpenXML sanitization {openxml_sanitized}."
                ),
            )
            output_size = output_file.stat().st_size
        if selected_variants:
            write_checkpoint(
                task_dir,
                "template-plan",
                {
                    "templateId": BUILTIN_TEMPLATE_ID,
                    "templateName": "电子科技大学",
                    "selectedVariants": selected_variants,
                    "roleMapping": role_mapping,
                },
            )
    else:
        openxml_sanitized = _sanitize_generated_pptx_openxml(output_file)
        openxml_errors = _validate_generated_pptx_openxml(output_file)
        if openxml_errors:
            raise RuntimeError("OpenXML validation failed after sanitization: " + "; ".join(openxml_errors[:8]))
        if any(openxml_sanitized.values()):
            append_log(task_dir, "info", f"OpenXML sanitization applied: {openxml_sanitized}.")
            output_size = output_file.stat().st_size
    _write_visual_pipeline_checkpoints(task_dir, project_dir, slide_count)
    preview_manifest = _write_svg_preview_manifest(task_dir, task_id, project_dir, slide_count)
    visual_pipeline_status, visual_fallback_reason, basic_fallback_detected = _visual_pipeline_status(
        project_dir,
        output_file,
        target_slide_count,
        settings,
        selected_variants,
        role_mapping,
    )
    fallback_reason = "; ".join(reason for reason in [repair_warning_reason, visual_fallback_reason] if reason) or None
    generation_mode = "paper-ppt-agent-rule-fallback" if basic_fallback_detected else "paper-ppt-agent"
    model_bridge_snapshot = _read_model_bridge_snapshot(task_dir)
    model_bridge_status = "failed" if basic_fallback_detected else model_bridge_snapshot.get("modelBridgeStatus") or "success"
    model_bridge_primary_status = "failed" if basic_fallback_detected else model_bridge_snapshot.get("modelBridgePrimaryStatus") or "success"
    model_bridge_fallback_status = (
        "failed" if basic_fallback_detected else model_bridge_snapshot.get("modelBridgeFallbackStatus") or "not_started"
    )
    model_bridge_error_summary = (
        fallback_reason
        if basic_fallback_detected or repair_warning_reason
        else model_bridge_snapshot.get("modelBridgeErrorSummary")
    )
    if visual_pipeline_status != "success":
        append_log(task_dir, "warn", fallback_reason or "Visual pipeline degraded; quality gates did not pass.")
    elif repair_warning_reason:
        append_log(task_dir, "warn", repair_warning_reason)
    write_checkpoint(
        task_dir,
        "pptx-exported",
        {
            "outputFileName": output_file.name,
            "outputFileSize": output_size,
            "slideCount": slide_count,
            "generator": "paper-ppt-agent",
            "generationMode": generation_mode,
            "visualPipelineStatus": visual_pipeline_status,
            "fallbackReason": fallback_reason,
            "modelBridgeStatus": model_bridge_status,
            "modelBridgePrimaryStatus": model_bridge_primary_status,
            "modelBridgeFallbackStatus": model_bridge_fallback_status,
            "modelBridgeErrorSummary": model_bridge_error_summary,
            "searchStatus": research_context["searchStatus"],
            "researchStatus": research_context["researchStatus"],
            "researchSourcesCount": research_context["researchSourcesCount"],
            "previewType": preview_manifest.get("type") if preview_manifest else None,
            "previewSlideCount": preview_manifest.get("slideCount") if preview_manifest else None,
            "selectedVariants": selected_variants,
            "roleMapping": role_mapping,
        },
    )
    append_log(task_dir, "info", f"PPTX generated. Size: {round(output_size / 1024)} KB.")
    return {
        "outputFileName": output_file.name,
        "outputFileSize": output_size,
        "slideCount": slide_count,
        "generationMode": generation_mode,
        "visualPipelineStatus": visual_pipeline_status,
        "fallbackReason": fallback_reason,
        "modelBridgeStatus": model_bridge_status,
        "modelBridgePrimaryModel": model_bridge_snapshot.get("modelBridgePrimaryModel") or "subrouter:gpt-5.4",
        "modelBridgePrimaryStatus": model_bridge_primary_status,
        "modelBridgeFallbackModel": model_bridge_snapshot.get("modelBridgeFallbackModel") or "moonshot:kimi-k2.5",
        "modelBridgeFallbackStatus": model_bridge_fallback_status,
        "modelBridgeErrorSummary": model_bridge_error_summary,
        "searchStatus": research_context["searchStatus"],
        "researchStatus": research_context["researchStatus"],
        "researchSourcesCount": research_context["researchSourcesCount"],
        "researchFallbackReason": research_context["researchFallbackReason"],
        "previewAvailable": bool(preview_manifest and preview_manifest.get("available")),
        "previewType": preview_manifest.get("type") if preview_manifest else None,
        "previewSlideCount": preview_manifest.get("slideCount") if preview_manifest else None,
        "previewManifestPath": str(safe_task_child(task_dir, "previews", "manifest.json")) if preview_manifest else None,
        "previewFallbackReason": (
            "预览仅供参考，最终以下载 PPTX 为准。"
            if _is_builtin_template_request(settings)
            else None
            if preview_manifest
            else "Real SVG preview is not available yet."
        ),
        "previewUpdatedAt": preview_manifest.get("generatedAt") if preview_manifest else None,
        "selectedVariants": selected_variants,
        "roleMapping": role_mapping,
    }
